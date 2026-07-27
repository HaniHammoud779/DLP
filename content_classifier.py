from transformers import pipeline
import re
import os
import json
import csv
import io
import zipfile
from html.parser import HTMLParser

from database.db import SessionLocal, OrganizationPolicy


try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None


try:
    from docx import Document
except Exception:
    Document = None


try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None


try:
    from pptx import Presentation
except Exception:
    Presentation = None


try:
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps
except Exception:
    Image = None
    ImageEnhance = None
    ImageFilter = None
    ImageOps = None


try:
    import pytesseract
except Exception:
    pytesseract = None


try:
    import fitz
except Exception:
    fitz = None


classifier = pipeline(
    "zero-shot-classification",
    model="facebook/bart-large-mnli"
)


MAX_EXTRACTED_CHARS = 20000

OCR_MIN_TEXT_CHARS = 25
OCR_MAX_PDF_PAGES = 20
OCR_IMAGE_SCALE = 2
OCR_TESSERACT_CONFIG = "--oem 3 --psm 6"

ORGANIZATION_POLICY_SENSITIVE_THRESHOLD = 0.72
ORGANIZATION_POLICY_MEDIUM_THRESHOLD = 0.64
ORGANIZATION_POLICY_SAFE_THRESHOLD = 0.60
ORGANIZATION_POLICY_MIN_MARGIN = 0.08

IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"
}


SECRET_SENSITIVE_THRESHOLD = 0.62
SECRET_MEDIUM_THRESHOLD = 0.55

BUSINESS_SENSITIVE_THRESHOLD = 0.88
BUSINESS_MEDIUM_THRESHOLD = 0.84
BUSINESS_MIN_MARGIN = 0.12
BUSINESS_MEDIUM_MIN_MARGIN = 0.18

GENERIC_SENSITIVE_THRESHOLD = 0.98
GENERIC_MEDIUM_THRESHOLD = 0.96

RESTAURANT_REFINEMENT_SENSITIVE_THRESHOLD = 0.86
RESTAURANT_REFINEMENT_SAFE_THRESHOLD = 0.70
RESTAURANT_REFINEMENT_MIN_MARGIN = 0.08

HEALTHCARE_REFINEMENT_SENSITIVE_THRESHOLD = 0.82
HEALTHCARE_REFINEMENT_MEDIUM_THRESHOLD = 0.72
HEALTHCARE_REFINEMENT_SENSITIVE_MARGIN = 0.08
HEALTHCARE_REFINEMENT_MEDIUM_MARGIN = 0.05


GENERIC_SENSITIVE_LABEL = (
    "This text contains confidential, private, restricted, or sensitive information."
)

GENERIC_SAFE_LABEL = (
    "This text is normal, public, operational, harmless, or safe."
)

GENERIC_LABELS = [
    GENERIC_SENSITIVE_LABEL,
    GENERIC_SAFE_LABEL
]


SECRET_EXPOSED_LABEL = (
    "This text exposes actual usable sensitive data that could be copied and misused, "
    "such as a real password, login credential, API key, access token, secret key, "
    "private key, payment card details, CVV, expiry date, bank account data, "
    "financial account data, or private customer account data."
)

SECRET_PLACEHOLDER_LABEL = (
    "This text contains only examples, placeholders, templates, dummy values, masked "
    "values, or instructions showing where a password, token, card number, key, or "
    "secret would be placed. It does not reveal a real usable secret."
)

SECRET_POLICY_LABEL = (
    "This text is only a policy, warning, recommendation, explanation, training text, "
    "or general discussion about passwords, tokens, cards, payments, banking, or "
    "sensitive data. It does not reveal actual usable sensitive data."
)

SECRET_PUBLIC_IDENTIFIER_LABEL = (
    "This text contains public or operational identifiers such as order identifiers, "
    "tracking identifiers, ticket identifiers, booking identifiers, invoice identifiers, "
    "or reference identifiers. It does not reveal payment-card details, credentials, "
    "tokens, or private account data."
)

SECRET_HARMLESS_LABEL = (
    "This text is a normal harmless message or ordinary communication. It does not "
    "contain real usable credentials, tokens, payment-card details, financial account "
    "data, private account data, or exposed secrets."
)

SECRET_REFINEMENT_LABELS = [
    SECRET_EXPOSED_LABEL,
    SECRET_PLACEHOLDER_LABEL,
    SECRET_POLICY_LABEL,
    SECRET_PUBLIC_IDENTIFIER_LABEL,
    SECRET_HARMLESS_LABEL
]


HEALTHCARE_LABEL = (
    "This text contains private healthcare information, such as a patient record, "
    "patient identity, diagnosis, prescription, treatment plan, laboratory result, "
    "medical test result, medical history, or clinical information about a person."
)

RESTAURANT_LABEL = (
    "This text exposes restaurant business-sensitive food production information, "
    "such as an internal recipe, sauce formula, ingredient quantities, seasoning "
    "formula, preparation timing, storage timing, or repeatable method used to "
    "produce a restaurant menu item."
)

SOFTWARE_LABEL = (
    "This text contains confidential software company intellectual property, such as "
    "private source code, proprietary algorithm logic, internal architecture, internal "
    "system design, unreleased technical plan, or implementation details not intended "
    "for public release."
)

CLOTHING_LABEL = (
    "This text contains confidential clothing retail business information, such as "
    "private customer records, supplier agreements, wholesale prices, purchase costs, "
    "discount negotiations, inventory information, internal pricing strategies, "
    "unreleased fashion collections, or confidential retail operations."
)

BUSINESS_SAFE_LABEL = (
    "This text is normal, harmless, public, personal, casual, or operational. It does "
    "not reveal private healthcare data, confidential business data, proprietary "
    "software information, or internal restaurant formula details."
)

BUSINESS_SEMANTIC_LABELS = [
    HEALTHCARE_LABEL,
    RESTAURANT_LABEL,
    SOFTWARE_LABEL,
    CLOTHING_LABEL,
    BUSINESS_SAFE_LABEL
]


HEALTHCARE_REFINEMENT_SENSITIVE_LABEL = (
    "This text clearly exposes concrete private healthcare information about a person, "
    "such as a patient identity together with a diagnosis, prescription, treatment, "
    "laboratory result, medical test result, medical history, clinical notes, or "
    "specific medical condition."
)

HEALTHCARE_REFINEMENT_GENERIC_LABEL = (
    "This text only mentions healthcare, a medical file, a doctor, a hospital, "
    "or a health-related topic in a general way. It does not reveal concrete "
    "private patient information, diagnosis, prescription, treatment, laboratory "
    "result, medical test result, medical history, or clinical details."
)

HEALTHCARE_REFINEMENT_OPERATIONAL_LABEL = (
    "This text is a normal operational or casual message related to healthcare, "
    "such as saying that a medical file exists, asking about an appointment, "
    "or referring to a doctor, without exposing private clinical details."
)

HEALTHCARE_REFINEMENT_LABELS = [
    HEALTHCARE_REFINEMENT_SENSITIVE_LABEL,
    HEALTHCARE_REFINEMENT_GENERIC_LABEL,
    HEALTHCARE_REFINEMENT_OPERATIONAL_LABEL
]


RESTAURANT_REFINEMENT_SENSITIVE_LABEL = (
    "This text clearly reveals enough concrete restaurant food-production information "
    "to reproduce a menu item or internal recipe. It includes concrete formula details "
    "such as quantities, measurements, preparation steps, mixing instructions, timing, "
    "storage conditions, or a complete repeatable method."
)

RESTAURANT_REFINEMENT_PUBLIC_LABEL = (
    "This text only mentions, describes, praises, reviews, or discusses food, a recipe, "
    "a sauce, a burger, a taste, or a menu item. It does not reveal enough concrete "
    "formula details to reproduce the item."
)

RESTAURANT_REFINEMENT_OPERATIONAL_LABEL = (
    "This text is normal restaurant operational information, such as opening hours, "
    "customer service, daily activity, public menu availability, or staff scheduling. "
    "It does not reveal internal food-production formula details."
)

RESTAURANT_REFINEMENT_LABELS = [
    RESTAURANT_REFINEMENT_SENSITIVE_LABEL,
    RESTAURANT_REFINEMENT_PUBLIC_LABEL,
    RESTAURANT_REFINEMENT_OPERATIONAL_LABEL
]


TEXT_EXTENSIONS = {
    ".txt", ".log", ".csv", ".json", ".xml", ".html", ".htm", ".md",
    ".py", ".java", ".cs", ".js", ".php", ".sql", ".env", ".conf",
    ".ini", ".yml", ".yaml", ".sh", ".bat", ".ps1", ".config"
}


class SimpleHTMLTextExtractor(HTMLParser):

    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, data):
        if data and data.strip():
            self.parts.append(data.strip())

    def get_text(self):
        return "\n".join(self.parts)


def limit_text(text):

    if not text:
        return ""

    if len(text) > MAX_EXTRACTED_CHARS:
        return text[:MAX_EXTRACTED_CHARS]

    return text


def read_text_file(file_path):

    try:
        with open(file_path, "r", errors="ignore", encoding="utf-8") as f:
            return f.read()
    except Exception:
        try:
            with open(file_path, "r", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""


def extract_html_text(file_path):

    raw_text = read_text_file(file_path)

    if not raw_text:
        return ""

    parser = SimpleHTMLTextExtractor()

    try:
        parser.feed(raw_text)
        return parser.get_text()
    except Exception:
        return raw_text



def ocr_is_available():

    return (
        Image is not None
        and pytesseract is not None
    )


def format_preview_section(title, text):

    clean_text = str(text or "").strip()

    if not clean_text:
        return ""

    return (
        "=" * 70
        + "\n"
        + title
        + "\n"
        + "=" * 70
        + "\n"
        + clean_text
    )


def combine_preview_sections(sections):

    clean_sections = [
        section.strip()
        for section in sections
        if section and section.strip()
    ]

    if not clean_sections:
        return ""

    return "\n\n".join(clean_sections)


def preprocess_image_for_ocr(image):

    if Image is None:
        return image

    try:
        prepared = image.convert("RGB")
    except Exception:
        prepared = image

    try:
        width, height = prepared.size

        if OCR_IMAGE_SCALE > 1:
            prepared = prepared.resize(
                (
                    max(1, width * OCR_IMAGE_SCALE),
                    max(1, height * OCR_IMAGE_SCALE)
                ),
                Image.Resampling.LANCZOS
            )
    except Exception:
        pass

    try:
        if ImageOps is not None:
            prepared = ImageOps.grayscale(prepared)
            prepared = ImageOps.autocontrast(prepared)
    except Exception:
        pass

    try:
        if ImageEnhance is not None:
            prepared = ImageEnhance.Contrast(
                prepared
            ).enhance(1.8)
    except Exception:
        pass

    try:
        if ImageFilter is not None:
            prepared = prepared.filter(
                ImageFilter.SHARPEN
            )
    except Exception:
        pass

    return prepared


def extract_ocr_text_from_pillow_image(image):

    if not ocr_is_available():
        return ""

    try:
        prepared_image = preprocess_image_for_ocr(
            image
        )

        extracted_text = pytesseract.image_to_string(
            prepared_image,
            config=OCR_TESSERACT_CONFIG
        )

        return str(extracted_text or "").strip()

    except Exception:
        return ""


def extract_image_ocr_text(file_path):

    if not ocr_is_available():
        return ""

    try:
        with Image.open(file_path) as image:
            return extract_ocr_text_from_pillow_image(
                image
            )

    except Exception:
        return ""


def extract_pdf_page_normal_texts(file_path):

    page_texts = []

    if PdfReader is None:
        return page_texts

    try:
        reader = PdfReader(file_path)

        for page in reader.pages:
            try:
                page_text = page.extract_text()
            except Exception:
                page_text = ""

            page_texts.append(
                str(page_text or "").strip()
            )

            if len(
                "\n".join(page_texts)
            ) >= MAX_EXTRACTED_CHARS:
                break

    except Exception:
        return []

    return page_texts


def extract_pdf_embedded_image_texts(page, document):

    image_texts = []

    if (
        not ocr_is_available()
        or fitz is None
    ):
        return image_texts

    try:
        page_images = page.get_images(
            full=True
        )
    except Exception:
        return image_texts

    seen_xrefs = set()

    for image_number, image_info in enumerate(
        page_images,
        start=1
    ):
        try:
            xref = image_info[0]

            if xref in seen_xrefs:
                continue

            seen_xrefs.add(xref)

            image_data = document.extract_image(
                xref
            )

            image_bytes = image_data.get(
                "image"
            )

            if not image_bytes:
                continue

            with Image.open(
                io.BytesIO(image_bytes)
            ) as image:
                image_text = (
                    extract_ocr_text_from_pillow_image(
                        image
                    )
                )

            if image_text:
                image_texts.append(
                    (
                        image_number,
                        image_text
                    )
                )

        except Exception:
            continue

    return image_texts


def extract_pdf_scanned_page_text(page):

    if (
        not ocr_is_available()
        or fitz is None
    ):
        return ""

    try:
        matrix = fitz.Matrix(
            2.0,
            2.0
        )

        pixmap = page.get_pixmap(
            matrix=matrix,
            alpha=False
        )

        image = Image.frombytes(
            "RGB",
            [
                pixmap.width,
                pixmap.height
            ],
            pixmap.samples
        )

        return extract_ocr_text_from_pillow_image(
            image
        )

    except Exception:
        return ""


def extract_pdf_content_with_sources(file_path):

    normal_page_texts = extract_pdf_page_normal_texts(
        file_path
    )

    classification_parts = []
    preview_sections = []

    fitz_document = None

    if fitz is not None:
        try:
            fitz_document = fitz.open(
                file_path
            )
        except Exception:
            fitz_document = None

    total_pages = len(
        normal_page_texts
    )

    if fitz_document is not None:
        total_pages = max(
            total_pages,
            min(
                fitz_document.page_count,
                OCR_MAX_PDF_PAGES
            )
        )

    total_pages = min(
        total_pages,
        OCR_MAX_PDF_PAGES
    )

    for page_index in range(
        total_pages
    ):
        page_number = page_index + 1

        normal_text = ""

        if page_index < len(
            normal_page_texts
        ):
            normal_text = normal_page_texts[
                page_index
            ].strip()

        if normal_text:
            classification_parts.append(
                normal_text
            )

            preview_sections.append(
                format_preview_section(
                    (
                        "DOCUMENT TEXT "
                        f"— PDF PAGE {page_number}"
                    ),
                    normal_text
                )
            )

        embedded_image_texts = []

        if (
            fitz_document is not None
            and page_index
            < fitz_document.page_count
        ):
            try:
                page = fitz_document.load_page(
                    page_index
                )

                embedded_image_texts = (
                    extract_pdf_embedded_image_texts(
                        page,
                        fitz_document
                    )
                )

                for (
                    image_number,
                    image_text
                ) in embedded_image_texts:
                    classification_parts.append(
                        image_text
                    )

                    preview_sections.append(
                        format_preview_section(
                            (
                                "OCR TEXT FROM EMBEDDED IMAGE "
                                f"{image_number} — PDF PAGE "
                                f"{page_number}"
                            ),
                            image_text
                        )
                    )

                if (
                    not embedded_image_texts
                    and len(
                        normal_text
                    )
                    < OCR_MIN_TEXT_CHARS
                ):
                    scanned_page_text = (
                        extract_pdf_scanned_page_text(
                            page
                        )
                    )

                    if scanned_page_text:
                        classification_parts.append(
                            scanned_page_text
                        )

                        preview_sections.append(
                            format_preview_section(
                                (
                                    "OCR TEXT FROM SCANNED "
                                    f"PDF PAGE {page_number}"
                                ),
                                scanned_page_text
                            )
                        )

            except Exception:
                pass

        if len(
            "\n".join(
                classification_parts
            )
        ) >= MAX_EXTRACTED_CHARS:
            break

    if fitz_document is not None:
        try:
            fitz_document.close()
        except Exception:
            pass

    classification_text = limit_text(
        "\n".join(
            classification_parts
        )
    )

    preview_text = limit_text(
        combine_preview_sections(
            preview_sections
        )
    )

    return (
        classification_text,
        preview_text
    )


def extract_pdf_text(file_path):

    classification_text, _ = (
        extract_pdf_content_with_sources(
            file_path
        )
    )

    return classification_text


def extract_docx_normal_text(file_path):

    if Document is None:
        return ""

    extracted_parts = []

    try:
        document = Document(file_path)

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()

            if text:
                extracted_parts.append(
                    text
                )

        for table in document.tables:
            for row in table.rows:
                row_values = []

                for cell in row.cells:
                    cell_text = cell.text.strip()

                    if cell_text:
                        row_values.append(
                            cell_text
                        )

                if row_values:
                    extracted_parts.append(
                        " | ".join(
                            row_values
                        )
                    )

    except Exception:
        return ""

    return "\n".join(
        extracted_parts
    )


def extract_docx_embedded_image_texts(file_path):

    image_texts = []

    if not ocr_is_available():
        return image_texts

    try:
        with zipfile.ZipFile(
            file_path,
            "r"
        ) as docx_archive:
            media_files = sorted(
                name
                for name in docx_archive.namelist()
                if name.startswith(
                    "word/media/"
                )
                and not name.endswith(
                    "/"
                )
            )

            for image_number, media_name in enumerate(
                media_files,
                start=1
            ):
                try:
                    image_bytes = docx_archive.read(
                        media_name
                    )

                    with Image.open(
                        io.BytesIO(
                            image_bytes
                        )
                    ) as image:
                        image_text = (
                            extract_ocr_text_from_pillow_image(
                                image
                            )
                        )

                    if image_text:
                        image_texts.append(
                            (
                                image_number,
                                os.path.basename(
                                    media_name
                                ),
                                image_text
                            )
                        )

                except Exception:
                    continue

    except Exception:
        return []

    return image_texts


def extract_docx_content_with_sources(file_path):

    normal_text = extract_docx_normal_text(
        file_path
    ).strip()

    embedded_image_texts = (
        extract_docx_embedded_image_texts(
            file_path
        )
    )

    classification_parts = []
    preview_sections = []

    if normal_text:
        classification_parts.append(
            normal_text
        )

        preview_sections.append(
            format_preview_section(
                "DOCUMENT TEXT — WORD FILE",
                normal_text
            )
        )

    for (
        image_number,
        image_name,
        image_text
    ) in embedded_image_texts:
        classification_parts.append(
            image_text
        )

        preview_sections.append(
            format_preview_section(
                (
                    "OCR TEXT FROM EMBEDDED IMAGE "
                    f"{image_number} — {image_name}"
                ),
                image_text
            )
        )

    return (
        limit_text(
            "\n".join(
                classification_parts
            )
        ),
        limit_text(
            combine_preview_sections(
                preview_sections
            )
        )
    )


def extract_docx_text(file_path):

    classification_text, _ = (
        extract_docx_content_with_sources(
            file_path
        )
    )

    return classification_text


def extract_office_embedded_image_texts(
    file_path,
    media_prefix,
):

    image_texts = []

    if not ocr_is_available():
        return image_texts

    try:
        with zipfile.ZipFile(
            file_path,
            "r"
        ) as office_archive:
            media_files = sorted(
                name
                for name in office_archive.namelist()
                if name.startswith(
                    media_prefix
                )
                and not name.endswith(
                    "/"
                )
            )

            for image_number, media_name in enumerate(
                media_files,
                start=1
            ):
                try:
                    image_bytes = office_archive.read(
                        media_name
                    )

                    with Image.open(
                        io.BytesIO(
                            image_bytes
                        )
                    ) as image:
                        image_text = (
                            extract_ocr_text_from_pillow_image(
                                image
                            )
                        )

                    if image_text:
                        image_texts.append(
                            (
                                image_number,
                                os.path.basename(
                                    media_name
                                ),
                                image_text
                            )
                        )

                except Exception:
                    continue

    except Exception:
        return []

    return image_texts


def extract_xlsx_normal_text(file_path):

    if load_workbook is None:
        return ""

    extracted_parts = []

    try:
        workbook = load_workbook(
            filename=file_path,
            read_only=True,
            data_only=True
        )

        for sheet in workbook.worksheets:
            extracted_parts.append(
                f"Sheet: {sheet.title}"
            )

            for row in sheet.iter_rows(
                values_only=True
            ):
                values = []

                for cell in row:
                    if cell is not None:
                        values.append(
                            str(cell)
                        )

                if values:
                    extracted_parts.append(
                        " | ".join(
                            values
                        )
                    )

                if len(
                    "\n".join(
                        extracted_parts
                    )
                ) >= MAX_EXTRACTED_CHARS:
                    break

            if len(
                "\n".join(
                    extracted_parts
                )
            ) >= MAX_EXTRACTED_CHARS:
                break

        workbook.close()

    except Exception:
        return ""

    return "\n".join(
        extracted_parts
    )


def extract_xlsx_embedded_image_texts(file_path):

    return extract_office_embedded_image_texts(
        file_path,
        "xl/media/"
    )


def extract_xlsx_content_with_sources(file_path):

    normal_text = extract_xlsx_normal_text(
        file_path
    ).strip()

    embedded_image_texts = (
        extract_xlsx_embedded_image_texts(
            file_path
        )
    )

    classification_parts = []
    preview_sections = []

    if normal_text:
        classification_parts.append(
            normal_text
        )

        preview_sections.append(
            format_preview_section(
                "DOCUMENT TEXT — EXCEL FILE",
                normal_text
            )
        )

    for (
        image_number,
        image_name,
        image_text
    ) in embedded_image_texts:
        classification_parts.append(
            image_text
        )

        preview_sections.append(
            format_preview_section(
                (
                    "OCR TEXT FROM EMBEDDED EXCEL IMAGE "
                    f"{image_number} — {image_name}"
                ),
                image_text
            )
        )

    return (
        limit_text(
            "\n".join(
                classification_parts
            )
        ),
        limit_text(
            combine_preview_sections(
                preview_sections
            )
        )
    )


def extract_xlsx_text(file_path):

    classification_text, _ = (
        extract_xlsx_content_with_sources(
            file_path
        )
    )

    return classification_text


def extract_pptx_normal_text(file_path):

    if Presentation is None:
        return ""

    extracted_parts = []

    try:
        presentation = Presentation(
            file_path
        )

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1
        ):
            extracted_parts.append(
                f"Slide {slide_number}"
            )

            for shape in slide.shapes:
                if hasattr(
                    shape,
                    "text"
                ):
                    text = shape.text.strip()

                    if text:
                        extracted_parts.append(
                            text
                        )

            if len(
                "\n".join(
                    extracted_parts
                )
            ) >= MAX_EXTRACTED_CHARS:
                break

    except Exception:
        return ""

    return "\n".join(
        extracted_parts
    )


def extract_pptx_embedded_image_texts(file_path):

    return extract_office_embedded_image_texts(
        file_path,
        "ppt/media/"
    )


def extract_pptx_content_with_sources(file_path):

    normal_text = extract_pptx_normal_text(
        file_path
    ).strip()

    embedded_image_texts = (
        extract_pptx_embedded_image_texts(
            file_path
        )
    )

    classification_parts = []
    preview_sections = []

    if normal_text:
        classification_parts.append(
            normal_text
        )

        preview_sections.append(
            format_preview_section(
                "DOCUMENT TEXT — POWERPOINT FILE",
                normal_text
            )
        )

    for (
        image_number,
        image_name,
        image_text
    ) in embedded_image_texts:
        classification_parts.append(
            image_text
        )

        preview_sections.append(
            format_preview_section(
                (
                    "OCR TEXT FROM EMBEDDED POWERPOINT IMAGE "
                    f"{image_number} — {image_name}"
                ),
                image_text
            )
        )

    return (
        limit_text(
            "\n".join(
                classification_parts
            )
        ),
        limit_text(
            combine_preview_sections(
                preview_sections
            )
        )
    )


def extract_pptx_text(file_path):

    classification_text, _ = (
        extract_pptx_content_with_sources(
            file_path
        )
    )

    return classification_text


def extract_csv_text(file_path):

    extracted_parts = []

    try:
        with open(file_path, "r", errors="ignore", encoding="utf-8") as f:
            reader = csv.reader(f)

            for row in reader:
                if row:
                    extracted_parts.append(" | ".join(row))

                if len("\n".join(extracted_parts)) >= MAX_EXTRACTED_CHARS:
                    break

    except Exception:
        return read_text_file(file_path)

    return "\n".join(extracted_parts)


def extract_json_text(file_path):

    raw_text = read_text_file(file_path)

    if not raw_text:
        return ""

    try:
        data = json.loads(raw_text)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return raw_text



def read_file_content_with_sources(file_path):

    extension = os.path.splitext(
        file_path
    )[1].lower()

    classification_text = ""
    preview_text = ""

    if extension in {
        ".html",
        ".htm"
    }:
        classification_text = extract_html_text(
            file_path
        )

        preview_text = format_preview_section(
            "DOCUMENT TEXT — HTML FILE",
            classification_text
        )

    elif extension == ".pdf":
        (
            classification_text,
            preview_text
        ) = extract_pdf_content_with_sources(
            file_path
        )

    elif extension in IMAGE_EXTENSIONS:
        classification_text = extract_image_ocr_text(
            file_path
        )

        preview_text = format_preview_section(
            (
                "OCR TEXT FROM IMAGE FILE "
                f"— {os.path.basename(file_path)}"
            ),
            classification_text
        )

    elif extension == ".docx":
        (
            classification_text,
            preview_text
        ) = extract_docx_content_with_sources(
            file_path
        )

    elif extension == ".xlsx":
        (
            classification_text,
            preview_text
        ) = extract_xlsx_content_with_sources(
            file_path
        )

    elif extension == ".pptx":
        (
            classification_text,
            preview_text
        ) = extract_pptx_content_with_sources(
            file_path
        )

    elif extension == ".csv":
        classification_text = extract_csv_text(
            file_path
        )

        preview_text = format_preview_section(
            "DOCUMENT TEXT — CSV FILE",
            classification_text
        )

    elif extension == ".json":
        classification_text = extract_json_text(
            file_path
        )

        preview_text = format_preview_section(
            "DOCUMENT TEXT — JSON FILE",
            classification_text
        )

    elif extension in TEXT_EXTENSIONS:
        classification_text = read_text_file(
            file_path
        )

        preview_text = format_preview_section(
            "DOCUMENT TEXT — TEXT FILE",
            classification_text
        )

    else:
        classification_text = read_text_file(
            file_path
        )

        preview_text = format_preview_section(
            "DOCUMENT TEXT",
            classification_text
        )

    classification_text = limit_text(
        classification_text
    )

    preview_text = limit_text(
        preview_text
    )

    return (
        classification_text,
        preview_text
    )


def read_file_content(file_path):

    classification_text, _ = (
        read_file_content_with_sources(
            file_path
        )
    )

    return classification_text


def read_file_preview(file_path):

    classification_text, preview_text = (
        read_file_content_with_sources(
            file_path
        )
    )

    if preview_text.strip():
        return preview_text

    return classification_text


def count_words(text):

    return len(re.findall(r"\b\w+\b", text or ""))


def has_reproducible_formula_detail(text):

    cleaned_text = text.strip()

    if not cleaned_text:
        return False

    lines = [
        line.strip()
        for line in cleaned_text.splitlines()
        if line.strip()
    ]

    number_pattern = re.compile(r"\b\d+(?:\.\d+)?\b")

    measurement_pattern = re.compile(
        r"\b\d+(?:\.\d+)?\s*"
        r"(ml|l|liter|liters|litre|litres|g|gram|grams|kg|mg|"
        r"tsp|tbsp|cup|cups|oz|lb|pound|pounds|minutes|minute|"
        r"hours|hour)\b",
        re.IGNORECASE
    )

    number_count = len(number_pattern.findall(cleaned_text))
    measurement_count = len(measurement_pattern.findall(cleaned_text))

    bullet_or_list_lines = 0

    for line in lines:
        if re.match(r"^\s*(-|\*|\d+\.|\d+\)|•)", line):
            bullet_or_list_lines += 1

    detail_score = 0

    if len(lines) >= 4:
        detail_score += 1

    if ":" in cleaned_text:
        detail_score += 1

    if number_count >= 3:
        detail_score += 1

    if measurement_count >= 2:
        detail_score += 2

    if bullet_or_list_lines >= 2:
        detail_score += 1

    return detail_score >= 3


def normalize_context_text(text):

    return " ".join((text or "").lower().split())


def has_explicit_healthcare_context(text):

    t = normalize_context_text(text)

    if not t:
        return False

    healthcare_patterns = [
        r"\bpatient\b",
        r"\bmedical\b",
        r"\bhealthcare\b",
        r"\bhealth care\b",
        r"\bhospital\b",
        r"\bclinic\b",
        r"\bclinical\b",
        r"\bdoctor\b",
        r"\bphysician\b",
        r"\bnurse\b",
        r"\bdiagnosis\b",
        r"\bdiagnosed\b",
        r"\bprescription\b",
        r"\bprescribed\b",
        r"\btreatment\b",
        r"\blaboratory\b",
        r"\blab result\b",
        r"\bmedical record\b",
        r"\bmedical history\b",
        r"\bmedical test\b"
    ]

    for pattern in healthcare_patterns:
        if re.search(pattern, t, re.IGNORECASE):
            return True

    return False


def build_safe_content_reason(text):

    t = normalize_context_text(text)

    if not t:
        return (
            "No readable sensitive information was detected. The content is "
            "treated as SAFE."
        )

    personal_introduction_patterns = [
        r"\bmy name is\b",
        r"\bi am \d{1,3} years? old\b",
        r"\bi['’]?m \d{1,3} years? old\b",
        r"\bmy age is \d{1,3}\b"
    ]

    if any(
        re.search(pattern, t, re.IGNORECASE)
        for pattern in personal_introduction_patterns
    ):
        return (
            "The content is an ordinary personal introduction containing general "
            "information. It does not expose credentials, "
            "payment-card data, private medical records, confidential business "
            "information, or other restricted data."
        )

    greeting_patterns = [
        r"\bhello\b",
        r"\bhi\b",
        r"\bgood morning\b",
        r"\bgood afternoon\b",
        r"\bgood evening\b"
    ]

    if count_words(t) <= 25 and any(
        re.search(pattern, t, re.IGNORECASE)
        for pattern in greeting_patterns
    ):
        return (
            "The content is a normal greeting or casual message. No credentials, "
            "financial data, private healthcare information, proprietary business "
            "information, or other restricted content was detected."
        )

    if has_explicit_healthcare_context(t):
        return (
            "The content mentions a healthcare-related topic, but it does not expose "
            "concrete private patient information such as a diagnosis, prescription, "
            "treatment, laboratory result, medical history, or clinical record."
        )

    return (
        "The content is normal, general, operational, personal, or public information. "
        "No usable credentials, payment-card data, private healthcare records, "
        "confidential business information, or other restricted data was confirmed."
    )


def has_concrete_clinical_detail(text):

    t = normalize_context_text(text)

    if not t:
        return False

    clinical_patterns = [
        r"\bdiagnosis\s*:",
        r"\bdiagnosed with\b",
        r"\bprescription\s*:",
        r"\bprescribed\b",
        r"\btreatment\s*:",
        r"\btreatment plan\b",
        r"\blab result\b",
        r"\blaboratory result\b",
        r"\bmedical test result\b",
        r"\bhba1c\b",
        r"\bblood glucose\b",
        r"\binsulin\b",
        r"\bdiabetes\b",
        r"\bhypertension\b",
        r"\bcancer\b",
        r"\bclinical notes\b",
        r"\bmedical history\b"
    ]

    for pattern in clinical_patterns:
        if re.search(pattern, t, re.IGNORECASE):
            return True

    return False


def has_healthcare_identity_context(text):

    t = normalize_context_text(text)

    if not t:
        return False

    # A general hospital reference must not be treated as patient identity.
    # This keeps public visiting hours, directions, and announcements SAFE.
    has_patient_context = re.search(
        r"\b(patient|patient file|medical file|medical record)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_identity_or_followup = (
        re.search(r"\bname\s*:", t, re.IGNORECASE) is not None
        or re.search(r"\bpatient name\s*:", t, re.IGNORECASE) is not None
        or re.search(r"\bpatient id\s*:", t, re.IGNORECASE) is not None
        or re.search(r"\bfollow[- ]?up\b", t, re.IGNORECASE) is not None
    )

    return has_patient_context and has_identity_or_followup


def has_incomplete_patient_file_context(text):

    t = normalize_context_text(text)

    if not t:
        return False

    has_patient_file_reference = re.search(
        r"\b(patient file|patient record|medical file|medical record)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_incomplete_or_review_context = re.search(
        r"\b(waiting for (?:security )?review|requires review|"
        r"not included|not shown|omitted|redacted|incomplete|partial)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_actual_private_detail = (
        has_concrete_clinical_detail(t)
        or re.search(
            r"\bpatient name\s*:\s*\S+",
            t,
            re.IGNORECASE
        ) is not None
        or re.search(
            r"\bpatient id\s*:\s*\S+",
            t,
            re.IGNORECASE
        ) is not None
    )

    return (
        has_patient_file_reference
        and has_incomplete_or_review_context
        and not has_actual_private_detail
    )


def has_public_healthcare_information(text):

    t = normalize_context_text(text)

    if not t:
        return False

    has_public_healthcare_context = re.search(
        r"\b(hospital|clinic|health center|medical center)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_public_operational_information = re.search(
        r"\b(visiting hours?|opening hours?|working hours?|office hours?|"
        r"contact information|phone number|address|directions|location|"
        r"parking|public announcement|visitor policy|appointment hours?)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_private_clinical_information = (
        has_healthcare_identity_context(t)
        or has_concrete_clinical_detail(t)
        or re.search(
            r"\b(patient id|patient name|diagnosis|prescription|"
            r"treatment plan|lab results?|medical history|clinical notes?)\b",
            t,
            re.IGNORECASE
        ) is not None
    )

    return (
        has_public_healthcare_context
        and has_public_operational_information
        and not has_private_clinical_information
    )


def has_confidential_warning_context(text):

    t = normalize_context_text(text)

    if not t:
        return False

    if re.search(
        r"\b(not confidential|not sensitive|public information|public document)\b",
        t,
        re.IGNORECASE
    ):
        return False

    has_confidential_marker = re.search(
        r"\b(confidential|sensitive|private|restricted|internal|do not share|dont share|not for public)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_document_or_business_context = re.search(
        r"\b(file|document|note|record|records|account|accounts|customer|finance|audit|team|company|business)\b",
        t,
        re.IGNORECASE
    ) is not None

    return has_confidential_marker and has_document_or_business_context


def has_medium_restaurant_context(text):

    t = normalize_context_text(text)

    if not t:
        return False

    has_food_process_context = re.search(
        r"\b(restaurant|sauce|recipe|ingredient|ingredients|kitchen|menu item|food production|mix|process)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_internal_or_incomplete_context = re.search(
        r"\b(internal|exact quantities|stored in the kitchen|do not share|dont share|not included|not shown)\b",
        t,
        re.IGNORECASE
    ) is not None

    return has_food_process_context and has_internal_or_incomplete_context


def has_medium_software_context(text):

    t = normalize_context_text(text)

    if not t:
        return False

    has_software_context = re.search(
        r"\b(system design|internal design|architecture|algorithm|scoring method|engine|source code|implementation)\b",
        t,
        re.IGNORECASE
    ) is not None

    has_internal_or_incomplete_context = re.search(
        r"\b(internal|private|confidential|not included|full algorithm is not included|proprietary)\b",
        t,
        re.IGNORECASE
    ) is not None

    return has_software_context and has_internal_or_incomplete_context


def has_medium_risk_context(text):

    if has_confidential_warning_context(text):
        return True

    if has_medium_restaurant_context(text):
        return True

    if has_medium_software_context(text):
        return True

    if has_healthcare_identity_context(text):
        return True

    return False


def is_email_address_shape(value):

    cleaned = value.strip().strip(".,!?;:'\"<>")

    email_pattern = re.compile(
        r"^[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}$"
    )

    return email_pattern.fullmatch(cleaned) is not None


def assignment_matches(text):

    pattern = re.compile(
        r"(?P<left>\b[A-Za-z_][A-Za-z0-9_\-]{1,50}\b)"
        r"(?P<separator>\s*(?:=|:|is)\s*)"
        r"(?P<quote>[\"']?)"
        r"(?P<value>[A-Za-z0-9@#$%^&*!_+=\-./<>*]{3,})"
        r"(?P=quote)",
        re.IGNORECASE
    )

    return list(pattern.finditer(text or ""))


def credential_exposure_matches(text):

    patterns = [
        re.compile(
            r"\b(?:my|the|this|user|admin|account|server|database|db|system|login)?\s*"
            r"(?P<kind>password|passcode|passwd|pwd)\s*"
            r"(?:=|:|is)\s*"
            r"(?P<value>[A-Za-z0-9@#$%^&*!_+=\-./<>*]{3,})",
            re.IGNORECASE
        ),
        re.compile(
            r"\b(?P<kind>api[_\- ]?key|secret[_\- ]?key|access[_\- ]?token|auth[_\- ]?token|private[_\- ]?key|token)\s*"
            r"(?:=|:|is)\s*"
            r"(?P<value>[A-Za-z0-9@#$%^&*!_+=\-./<>*]{6,})",
            re.IGNORECASE
        )
    ]

    matches = []

    for pattern in patterns:
        matches.extend(list(pattern.finditer(text or "")))

    return matches


def looks_like_readable_placeholder_shape(value):

    cleaned = value.strip().strip(".,!?;:'\"")

    if not cleaned:
        return False

    if is_email_address_shape(cleaned):
        return False

    if re.fullmatch(r"<[A-Za-z0-9_\- ]+>", cleaned):
        return True

    if re.fullmatch(r"\*{3,}", cleaned):
        return True

    if re.fullmatch(r"[A-Za-z_]+", cleaned) and "_" in cleaned:
        return True

    if re.fullmatch(r"[A-Z_]+", cleaned) and "_" in cleaned:
        return True

    return False


def looks_like_secret_value_shape(value):

    cleaned = value.strip().strip(".,!?;:'\"")

    if len(cleaned) < 6:
        return False

    if is_email_address_shape(cleaned):
        return False

    if looks_like_readable_placeholder_shape(cleaned):
        return False

    has_letter = re.search(r"[A-Za-z]", cleaned) is not None
    has_digit = re.search(r"\d", cleaned) is not None
    strong_special = re.search(r"[@#$%^&*!+=\-./]", cleaned) is not None

    if cleaned.isdigit() and len(cleaned) >= 6:
        return True

    if cleaned.isdigit():
        return False

    if has_letter and has_digit:
        return True

    if has_letter and strong_special and len(cleaned) >= 8:
        return True

    if has_digit and strong_special and len(cleaned) >= 8:
        return True

    return False


def looks_like_strong_credential_value(value, kind):

    cleaned = value.strip().strip(".,!?;:'\"")

    if not cleaned:
        return False

    if is_email_address_shape(cleaned):
        return False

    if looks_like_readable_placeholder_shape(cleaned):
        return False

    kind_lower = (kind or "").lower()

    has_letter = re.search(r"[A-Za-z]", cleaned) is not None
    has_digit = re.search(r"\d", cleaned) is not None
    has_special = re.search(r"[@#$%^&*!+=\-./]", cleaned) is not None

    if kind_lower in ["password", "passcode", "passwd", "pwd"]:
        return looks_like_secret_value_shape(cleaned)

    if "token" in kind_lower or "key" in kind_lower:

        if re.search(
            r"\b(test|demo|sample|example|dummy|placeholder|replace|your)\b",
            cleaned,
            re.IGNORECASE
        ):
            return False

        if len(cleaned) >= 20 and has_letter and (has_digit or has_special):
            return True

        if len(cleaned) >= 16 and has_letter and has_digit and has_special:
            return True

    return False


def has_credential_exposure_shape(text):

    for match in credential_exposure_matches(text):

        kind = match.group("kind")
        value = match.group("value").strip().strip(".,!?;:'\"")

        if not value:
            continue

        if is_email_address_shape(value):
            continue

        if looks_like_readable_placeholder_shape(value):
            continue

        if looks_like_strong_credential_value(value, kind):
            return True

    return False


def has_weak_credential_context(text):

    for match in credential_exposure_matches(text):

        kind = match.group("kind")
        value = match.group("value").strip().strip(".,!?;:'\"")

        if not value:
            continue

        if is_email_address_shape(value):
            continue

        if looks_like_readable_placeholder_shape(value):
            continue

        if looks_like_strong_credential_value(value, kind):
            continue

        kind_lower = (kind or "").lower()

        if "token" in kind_lower or "key" in kind_lower:
            has_letter = re.search(r"[A-Za-z]", value) is not None
            has_digit = re.search(r"\d", value) is not None

            if len(value) >= 6 and has_letter and has_digit:
                return True

    return False


def has_assignment_value_shape(text):

    for match in assignment_matches(text):

        left = match.group("left").strip().lower()
        value = match.group("value").strip().strip(".,!?;:'\"")

        if is_email_address_shape(value):
            continue

        if left in [
            "name",
            "patient",
            "diagnosis",
            "results",
            "result",
            "plan",
            "subject",
            "recipient"
        ]:
            continue

        if looks_like_secret_value_shape(value):
            return True

    return False


def has_long_token_like_shape(text):

    tokens = re.findall(r"\b[A-Za-z0-9@#$%^&*!_+=\-./]{16,}\b", text or "")

    for token in tokens:

        if is_email_address_shape(token):
            continue

        if looks_like_readable_placeholder_shape(token):
            continue

        has_letter = re.search(r"[A-Za-z]", token) is not None
        has_digit = re.search(r"\d", token) is not None
        strong_special = re.search(r"[@#$%^&*!+=\-./]", token) is not None

        if has_letter and (has_digit or strong_special):
            return True

    return False


def extract_card_number_candidates(text):

    if not text:
        return []

    grouped_pattern = re.compile(
        r"(?<!\d)(?:\d[ -]?){13,19}(?!\d)"
    )

    candidates = []

    for match in grouped_pattern.finditer(text):

        original_value = match.group(0).strip()

        digits_only = re.sub(
            r"\D",
            "",
            original_value
        )

        if 13 <= len(digits_only) <= 19:

            candidates.append(
                {
                    "original": original_value,
                    "digits": digits_only
                }
            )

    return candidates


def passes_luhn_check(number):

    digits = re.sub(
        r"\D",
        "",
        str(number or "")
    )

    if not 13 <= len(digits) <= 19:
        return False

    if len(set(digits)) == 1:
        return False

    total = 0
    reverse_digits = digits[::-1]

    for index, character in enumerate(
        reverse_digits
    ):

        digit = int(character)

        if index % 2 == 1:

            digit *= 2

            if digit > 9:
                digit -= 9

        total += digit

    return total % 10 == 0


def has_valid_payment_card_number(text):

    candidates = extract_card_number_candidates(
        text
    )

    for candidate in candidates:

        if passes_luhn_check(
            candidate["digits"]
        ):
            return True

    return False


def has_card_number_shape(text):

    return len(
        extract_card_number_candidates(
            text
        )
    ) > 0


def has_card_supporting_numeric_shape(text):

    if not text:
        return False

    expiry_shape = re.search(
        r"\b(0[1-9]|1[0-2])\s*/\s*(\d{2}|\d{4})\b",
        text
    ) is not None

    short_security_code_shape = re.search(
        r"\b\d{3,4}\b",
        text
    ) is not None

    return expiry_shape or short_security_code_shape


def has_payment_card_structure(text):

    if has_valid_payment_card_number(text):
        return True

    return (
        has_card_number_shape(text)
        and has_card_supporting_numeric_shape(text)
    )


def has_secret_structural_evidence(text):

    if has_credential_exposure_shape(text):
        return True

    if has_payment_card_structure(text):
        return True

    if has_assignment_value_shape(text):
        return True

    if has_long_token_like_shape(text):
        return True

    return False


def mask_word(value):

    return "*" * len(value)


def mask_sensitive_content(text):

    sensitive_snippets = []

    for line in text.splitlines():

        clean_line = line.strip()

        if not clean_line:
            continue

        matches = assignment_matches(clean_line)

        for match in matches:

            left = match.group("left")
            separator = match.group("separator")
            quote = match.group("quote")
            value = match.group("value")

            if is_email_address_shape(value):
                continue

            if looks_like_secret_value_shape(value):

                snippet = (
                    mask_word(left)
                    + separator
                    + quote
                    + mask_word(value)
                    + quote
                )

                sensitive_snippets.append(snippet)

    card_matches = re.findall(
        r"\b(?:\d[ -]?){13,19}\b",
        text
    )

    for card in card_matches:
        sensitive_snippets.append(mask_word(card))

    long_tokens = re.findall(
        r"\b[A-Za-z0-9@#$%^&*!_+=\-./]{16,}\b",
        text
    )

    for token in long_tokens:

        if is_email_address_shape(token):
            continue

        if looks_like_readable_placeholder_shape(token):
            continue

        if looks_like_secret_value_shape(token):
            sensitive_snippets.append(mask_word(token))

    if not sensitive_snippets:

        return (
            "Sensitive content detected semantically by AI. "
            "No direct secret-like value was extracted for display."
        )

    unique_snippets = []

    for snippet in sensitive_snippets:

        if snippet not in unique_snippets:
            unique_snippets.append(snippet)

    return "\n".join(unique_snippets)


def run_zero_shot(content, labels):

    try:
        result = classifier(
            content,
            labels,
            hypothesis_template="{}",
            truncation=True
        )

        return result

    except Exception as e:
        return {
            "labels": [labels[-1]],
            "scores": [0],
            "error": str(e)
        }


def run_secret_refinement_ai_analysis(content):

    result = run_zero_shot(content, SECRET_REFINEMENT_LABELS)

    labels = result["labels"]
    scores = result["scores"]

    best_label = labels[0]
    best_score = float(scores[0])

    second_score = 0

    if len(scores) > 1:
        second_score = float(scores[1])

    return {
        "best_label": best_label,
        "best_score": best_score,
        "second_score": second_score,
        "margin": best_score - second_score
    }


def run_generic_ai_analysis(content):

    result = run_zero_shot(content, GENERIC_LABELS)

    return {
        "best_label": result["labels"][0],
        "best_score": float(result["scores"][0])
    }


def run_business_ai_analysis(content):

    result = run_zero_shot(content, BUSINESS_SEMANTIC_LABELS)

    labels = result["labels"]
    scores = result["scores"]

    best_label = labels[0]
    best_score = float(scores[0])

    safe_score = 0

    for label, score in zip(labels, scores):
        if label == BUSINESS_SAFE_LABEL:
            safe_score = float(score)
            break

    return {
        "best_label": best_label,
        "best_score": best_score,
        "safe_score": safe_score,
        "margin_over_safe": best_score - safe_score
    }


def run_healthcare_refinement_ai_analysis(content):

    result = run_zero_shot(content, HEALTHCARE_REFINEMENT_LABELS)

    labels = result["labels"]
    scores = result["scores"]

    best_label = labels[0]
    best_score = float(scores[0])

    second_score = 0

    if len(scores) > 1:
        second_score = float(scores[1])

    return {
        "best_label": best_label,
        "best_score": best_score,
        "second_score": second_score,
        "margin": best_score - second_score
    }


def run_restaurant_refinement_ai_analysis(content):

    result = run_zero_shot(content, RESTAURANT_REFINEMENT_LABELS)

    labels = result["labels"]
    scores = result["scores"]

    best_label = labels[0]
    best_score = float(scores[0])

    second_score = 0

    if len(scores) > 1:
        second_score = float(scores[1])

    return {
        "best_label": best_label,
        "best_score": best_score,
        "second_score": second_score,
        "margin": best_score - second_score
    }


def get_business_reason(best_business_label):

    if best_business_label == HEALTHCARE_LABEL:
        return (
            "AI semantic analysis detected private healthcare information, "
            "such as patient records, diagnoses, prescriptions, lab results, "
            "treatment details, or medical history."
        )

    if best_business_label == RESTAURANT_LABEL:
        return (
            "AI semantic analysis detected restaurant business-sensitive recipe "
            "or food production information."
        )

    if best_business_label == SOFTWARE_LABEL:
        return (
            "AI semantic analysis detected software-company intellectual property, "
            "such as proprietary algorithms, private source code, internal system "
            "design, architecture details, or unreleased technical plans."
        )

    if best_business_label == CLOTHING_LABEL:
        return (
            "AI semantic analysis detected confidential clothing retail information, "
            "such as private customer records, supplier details, wholesale prices, "
            "purchase costs, discounts, inventory, or unreleased fashion plans."
        )

    return ""



def get_active_organization_policy():

    session = SessionLocal()

    try:
        policy = session.query(
            OrganizationPolicy
        ).filter(
            OrganizationPolicy.is_active.is_(True)
        ).order_by(
            OrganizationPolicy.updated_at.desc(),
            OrganizationPolicy.id.desc()
        ).first()

        if not policy:
            return {
                "organization_name": "",
                "policy_prompt": "",
                "policy_version": 0,
                "is_active": False
            }

        return {
            "organization_name": str(policy.organization_name or "").strip(),
            "policy_prompt": str(policy.policy_prompt or "").strip(),
            "policy_version": int(policy.policy_version or 1),
            "is_active": bool(policy.is_active)
        }

    except Exception as error:
        print(
            "[ORGANIZATION POLICY WARNING] "
            f"Could not load policy: {error}"
        )

        return {
            "organization_name": "",
            "policy_prompt": "",
            "policy_version": 0,
            "is_active": False
        }

    finally:
        session.close()


def parse_organization_policy_sections(policy_prompt):

    prompt = str(policy_prompt or "").replace("\r", "\n").strip()

    sections = {
        "scope": [],
        "sensitive": [],
        "medium": [],
        "safe": [],
        "outside_scope": []
    }

    if not prompt:
        return sections

    heading_map = {
        "ORGANIZATION SCOPE": "scope",
        "ORGANISATION SCOPE": "scope",
        "BUSINESS SCOPE": "scope",
        "SCOPE": "scope",
        "SENSITIVE": "sensitive",
        "MEDIUM": "medium",
        "SAFE": "safe",
        "OUTSIDE SCOPE": "outside_scope",
        "OUT OF SCOPE": "outside_scope"
    }

    current_section = None

    for raw_line in prompt.splitlines():

        line = str(raw_line or "").strip()

        if not line:
            continue

        heading_match = re.match(
            r"^([A-Za-z ]{3,40})\s*:\s*(.*)$",
            line
        )

        if heading_match:

            possible_heading = re.sub(
                r"[^A-Za-z ]",
                "",
                heading_match.group(1)
            ).strip().upper()

            if possible_heading in heading_map:

                current_section = heading_map[
                    possible_heading
                ]

                remainder = heading_match.group(2).strip()

                if remainder:
                    sections[current_section].append(
                        remainder
                    )

                continue

        heading_candidate = re.sub(
            r"[^A-Za-z ]",
            "",
            line
        ).strip().upper()

        if heading_candidate in heading_map:
            current_section = heading_map[
                heading_candidate
            ]
            continue

        if current_section is not None:

            cleaned_line = re.sub(
                r"^\s*(?:[-*•]|\d+[.)])\s*",
                "",
                line
            ).strip()

            if cleaned_line:
                sections[current_section].append(
                    cleaned_line
                )

    return sections


def join_policy_section(items):

    clean_items = [
        str(item or "").strip()
        for item in items
        if str(item or "").strip()
    ]

    return " ".join(clean_items)



ORGANIZATION_DOMAIN_LABELS = {
    "HEALTHCARE": (
        "This content belongs to healthcare, hospitals, clinics, medicine, "
        "clinical care, patients, medical records, insurance, or patient services."
    ),
    "FOOD_HOSPITALITY": (
        "This content belongs to restaurants, food service, kitchens, recipes, "
        "hospitality, catering, beverages, hotels, or food production."
    ),
    "SOFTWARE_TECHNOLOGY": (
        "This content belongs to software, source code, algorithms, cybersecurity, "
        "technical architecture, information technology, digital platforms, or systems."
    ),
    "RETAIL_COMMERCE": (
        "This content belongs to retail, wholesale, clothing, textiles, fashion, "
        "merchandise, product pricing, suppliers, sales, or commercial distribution."
    ),
    "FINANCE_BANKING": (
        "This content belongs to banking, finance, accounting, investments, loans, "
        "payments, insurance finance, fraud analytics, or financial services."
    ),
    "EDUCATION": (
        "This content belongs to schools, universities, courses, students, teaching, "
        "academic administration, research education, or training."
    ),
    "CONSTRUCTION_ENGINEERING": (
        "This content belongs to construction, civil engineering, architecture, "
        "building projects, contracting, infrastructure, or real-estate development."
    ),
    "MANUFACTURING_INDUSTRIAL": (
        "This content belongs to manufacturing, factories, industrial production, "
        "machinery, product formulas, quality control, or supply-chain operations."
    ),
    "LEGAL_GOVERNMENT": (
        "This content belongs to legal services, courts, law firms, government, "
        "public administration, regulation, compliance, or official records."
    ),
    "LOGISTICS_TRANSPORT": (
        "This content belongs to logistics, shipping, transport, warehousing, "
        "delivery, aviation, ports, fleets, or freight operations."
    ),
    "MEDIA_COMMUNICATIONS": (
        "This content belongs to media, journalism, advertising, publishing, "
        "telecommunications, public relations, or communications."
    ),
    "ENERGY_UTILITIES": (
        "This content belongs to energy, electricity, oil, gas, utilities, "
        "renewable power, water services, or resource operations."
    ),
    "AGRICULTURE": (
        "This content belongs to agriculture, farming, crops, livestock, "
        "food cultivation, agricultural supply, or rural production."
    ),
    "GENERAL_OTHER": (
        "This content belongs to another professional or business domain not "
        "described by the other categories."
    )
}


def run_organization_domain_analysis(text):

    clean_text = str(text or "").strip()

    if not clean_text:
        return {
            "domain": "NONE",
            "score": 0,
            "margin": 0
        }

    label_list = list(
        ORGANIZATION_DOMAIN_LABELS.values()
    )

    try:
        result = classifier(
            clean_text,
            label_list,
            multi_label=False,
            hypothesis_template="{}",
            truncation=True
        )

        labels = result.get(
            "labels",
            []
        )

        scores = result.get(
            "scores",
            []
        )

        if not labels or not scores:
            return {
                "domain": "NONE",
                "score": 0,
                "margin": 0
            }

        best_label = labels[0]
        best_score = float(
            scores[0]
        )

        second_score = (
            float(
                scores[1]
            )
            if len(
                scores
            ) > 1
            else 0
        )

        reverse_map = {
            value: key
            for key, value
            in ORGANIZATION_DOMAIN_LABELS.items()
        }

        return {
            "domain": reverse_map.get(
                best_label,
                "NONE"
            ),
            "score": best_score,
            "margin": best_score - second_score
        }

    except Exception as error:
        print(
            "[ORGANIZATION POLICY WARNING] "
            f"Domain analysis failed: {error}"
        )

        return {
            "domain": "NONE",
            "score": 0,
            "margin": 0
        }


def run_direct_scope_alignment_analysis(
    content,
    organization_name,
    scope_text,
    outside_scope_text
):

    clean_content = str(
        content or ""
    ).strip()

    clean_scope = str(
        scope_text or ""
    ).strip()

    if not clean_content or not clean_scope:
        return {
            "classification": "NONE",
            "score": 0,
            "margin": 0
        }

    organization_label = (
        "This content directly concerns the active organization, "
        f"{organization_name or 'the organization'}, and its stated business "
        f"scope: {clean_scope}"
    )

    unrelated_label = (
        "This content belongs to a different external industry or business "
        "domain and is unrelated to the active organization's stated scope."
    )

    if outside_scope_text:
        unrelated_label += (
            " The organization policy says unrelated content is outside scope: "
            f"{outside_scope_text}"
        )

    try:
        result = classifier(
            clean_content,
            [
                organization_label,
                unrelated_label
            ],
            multi_label=False,
            hypothesis_template="{}",
            truncation=True
        )

        labels = result.get(
            "labels",
            []
        )

        scores = result.get(
            "scores",
            []
        )

        if not labels or not scores:
            return {
                "classification": "NONE",
                "score": 0,
                "margin": 0
            }

        best_label = labels[0]
        best_score = float(
            scores[0]
        )

        second_score = (
            float(
                scores[1]
            )
            if len(
                scores
            ) > 1
            else 0
        )

        return {
            "classification": (
                "IN_SCOPE"
                if best_label == organization_label
                else "OUT_OF_SCOPE"
            ),
            "score": best_score,
            "margin": best_score - second_score
        }

    except Exception as error:
        print(
            "[ORGANIZATION POLICY WARNING] "
            f"Direct scope analysis failed: {error}"
        )

        return {
            "classification": "NONE",
            "score": 0,
            "margin": 0
        }


def run_organization_policy_ai_analysis(
    content,
    policy_prompt,
    organization_name=""
):

    clean_content = str(content or "").strip()
    sections = parse_organization_policy_sections(
        policy_prompt
    )

    scope_text = join_policy_section(
        sections.get(
            "scope",
            []
        )
    )

    sensitive_text = join_policy_section(
        sections.get(
            "sensitive",
            []
        )
    )

    medium_text = join_policy_section(
        sections.get(
            "medium",
            []
        )
    )

    safe_text = join_policy_section(
        sections.get(
            "safe",
            []
        )
    )

    outside_scope_text = join_policy_section(
        sections.get(
            "outside_scope",
            []
        )
    )

    organization_domain_result = run_organization_domain_analysis(
        scope_text
    )

    content_domain_result = run_organization_domain_analysis(
        clean_content
    )

    organization_domain = organization_domain_result.get(
        "domain",
        "NONE"
    )

    content_domain = content_domain_result.get(
        "domain",
        "NONE"
    )

    # Scope alignment should consider the complete customer policy.
    # The scope section defines the organization, while sensitive and
    # medium sections define protected business information that belongs
    # to that organization.
    policy_scope_alignment_text = " ".join(
        [
            scope_text,
            sensitive_text,
            medium_text
        ]
    ).strip()

    direct_scope_result = run_direct_scope_alignment_analysis(
        clean_content,
        organization_name,
        policy_scope_alignment_text,
        outside_scope_text
    )

    direct_scope_classification = direct_scope_result.get(
        "classification",
        "NONE"
    )

    direct_scope_score = float(
        direct_scope_result.get(
            "score",
            0
        )
    )

    direct_scope_margin = float(
        direct_scope_result.get(
            "margin",
            0
        )
    )

    recognized_domains = set(
        ORGANIZATION_DOMAIN_LABELS.keys()
    ) - {
        "GENERAL_OTHER"
    }

    strong_domain_mismatch = (
        organization_domain in recognized_domains
        and content_domain in recognized_domains
        and organization_domain != content_domain
        and organization_domain_result.get(
            "score",
            0
        ) >= 0.42
        and content_domain_result.get(
            "score",
            0
        ) >= 0.42
    )

    # If the document domain matches the organization domain, do not allow
    # a generic direct scope classifier to incorrectly mark the content as
    # outside the organization. The detailed policy sections decide the
    # final SAFE/MEDIUM/SENSITIVE level afterwards.
    same_business_domain = (
        organization_domain in recognized_domains
        and content_domain in recognized_domains
        and organization_domain == content_domain
    )

    if same_business_domain:
        strong_direct_outside_scope = False

    policy_protection_result = run_zero_shot(
        clean_content,
        [
            (
                "This content matches the active organization's protected "
                "sensitive or medium information described in the policy: "
                f"{sensitive_text} {medium_text}"
            ),
            (
                "This content does not match the organization's protected "
                "information categories."
            )
        ]
    )

    policy_protection_score = float(
        policy_protection_result.get(
            "scores",
            [0]
        )[0]
    )

    strong_policy_match = (
        policy_protection_score >= 0.65
        and (
            organization_domain == content_domain
            or organization_domain == "NONE"
            or content_domain == "NONE"
        )
    )

    # A clear organization-domain match should prevent an incorrect
    # scope rejection. The customer's own policy should control content
    # inside its business domain.
    organization_policy_domain_match = (
        organization_domain != "NONE"
        and content_domain != "NONE"
        and organization_domain == content_domain
        and strong_policy_match
    )

    # Do not allow the generic scope classifier to reject content when the
    # organization and the document are already identified as the same business
    # domain. The organization policy should decide the risk level (SAFE/MEDIUM/
    # SENSITIVE) inside its own domain. This prevents restaurant documents such
    # as internal recipe notes from becoming OUT_OF_SCOPE because the text does
    # not contain the exact scope wording.
    strong_direct_outside_scope = (
        direct_scope_classification == "OUT_OF_SCOPE"
        and direct_scope_score >= 0.52
        and not same_business_domain
        and not organization_policy_domain_match
        and not strong_policy_match
    )

    # A document is outside scope when either:
    # 1) the broad AI business domains clearly differ, or
    # 2) the direct AI scope comparison says it belongs to another industry.
    #
    # Mandatory universal protections are evaluated separately later and
    # cannot be downgraded by this organization-scope decision.
    if (
        strong_domain_mismatch
        or strong_direct_outside_scope
    ):
        scope_score = max(
            float(
                organization_domain_result.get(
                    "score",
                    0
                )
            ),
            float(
                content_domain_result.get(
                    "score",
                    0
                )
            ),
            direct_scope_score
        )

        scope_margin = max(
            float(
                organization_domain_result.get(
                    "margin",
                    0
                )
            ),
            float(
                content_domain_result.get(
                    "margin",
                    0
                )
            ),
            direct_scope_margin
        )

        return {
            "classification": "OUT_OF_SCOPE",
            "best_score": scope_score,
            "margin": scope_margin,
            "scope_classification": "OUT_OF_SCOPE",
            "scope_score": scope_score,
            "scope_margin": scope_margin,
            "matched_section": (
                f"Active organization domain: {organization_domain}. "
                f"Content domain: {content_domain}. "
                f"Direct scope result: {direct_scope_classification}."
            )
        }

    if not clean_content or not any(
        [
            scope_text,
            sensitive_text,
            medium_text,
            safe_text,
            outside_scope_text
        ]
    ):
        return {
            "classification": "NONE",
            "best_score": 0,
            "margin": 0,
            "scope_classification": "NONE",
            "scope_score": 0,
            "scope_margin": 0,
            "matched_section": ""
        }

    scope_classification = "NONE"
    scope_score = 0.0
    scope_margin = 0.0

    if scope_text:

        in_scope_label = (
            "This content directly belongs to the active organization's "
            "business scope and activities described here: "
            f"{scope_text}"
        )

        outside_scope_label = (
            "This content belongs to another unrelated industry or domain "
            "and is outside the active organization's business scope."
        )

        if outside_scope_text:
            outside_scope_label += (
                " The policy gives these outside-scope descriptions or "
                f"examples: {outside_scope_text}"
            )

        try:
            scope_result = classifier(
                clean_content,
                [
                    in_scope_label,
                    outside_scope_label
                ],
                multi_label=False,
                hypothesis_template="{}",
                truncation=True
            )

            returned_labels = scope_result.get(
                "labels",
                []
            )

            returned_scores = scope_result.get(
                "scores",
                []
            )

            if returned_labels and returned_scores:

                top_label = returned_labels[0]
                scope_score = float(
                    returned_scores[0]
                )

                second_scope_score = (
                    float(
                        returned_scores[1]
                    )
                    if len(
                        returned_scores
                    ) > 1
                    else 0.0
                )

                scope_margin = (
                    scope_score
                    - second_scope_score
                )

                if top_label == in_scope_label:
                    scope_classification = "IN_SCOPE"
                else:
                    scope_classification = "OUT_OF_SCOPE"

        except Exception as error:
            print(
                "[ORGANIZATION POLICY WARNING] "
                f"Scope analysis failed: {error}"
            )

    if (
        scope_classification == "OUT_OF_SCOPE"
        and scope_score >= 0.60
        and scope_margin >= 0.08
    ):
        return {
            "classification": "OUT_OF_SCOPE",
            "best_score": scope_score,
            "margin": scope_margin,
            "scope_classification": scope_classification,
            "scope_score": scope_score,
            "scope_margin": scope_margin,
            "matched_section": outside_scope_text
        }

    labels = []
    label_to_classification = {}

    if sensitive_text:

        sensitive_label = (
            "According to the active organization policy, this content "
            "matches the SENSITIVE category: "
            f"{sensitive_text}"
        )

        labels.append(
            sensitive_label
        )

        label_to_classification[
            sensitive_label
        ] = "SENSITIVE"

    if medium_text:

        medium_label = (
            "According to the active organization policy, this content "
            "matches the MEDIUM category and requires review: "
            f"{medium_text}"
        )

        labels.append(
            medium_label
        )

        label_to_classification[
            medium_label
        ] = "MEDIUM"

    if safe_text:

        safe_label = (
            "According to the active organization policy, this content "
            "matches the SAFE or public category: "
            f"{safe_text}"
        )

        labels.append(
            safe_label
        )

        label_to_classification[
            safe_label
        ] = "SAFE"

    if not labels:
        return {
            "classification": "NONE",
            "best_score": 0,
            "margin": 0,
            "scope_classification": scope_classification,
            "scope_score": scope_score,
            "scope_margin": scope_margin,
            "matched_section": ""
        }

    try:
        result = classifier(
            clean_content,
            labels,
            multi_label=False,
            hypothesis_template="{}",
            truncation=True
        )

        returned_labels = result.get(
            "labels",
            []
        )

        returned_scores = result.get(
            "scores",
            []
        )

    except Exception as error:
        print(
            "[ORGANIZATION POLICY WARNING] "
            f"Risk analysis failed: {error}"
        )

        return {
            "classification": "NONE",
            "best_score": 0,
            "margin": 0,
            "scope_classification": scope_classification,
            "scope_score": scope_score,
            "scope_margin": scope_margin,
            "matched_section": ""
        }

    if not returned_labels or not returned_scores:
        return {
            "classification": "NONE",
            "best_score": 0,
            "margin": 0,
            "scope_classification": scope_classification,
            "scope_score": scope_score,
            "scope_margin": scope_margin,
            "matched_section": ""
        }

    best_label = returned_labels[0]
    best_score = float(
        returned_scores[0]
    )

    second_score = (
        float(
            returned_scores[1]
        )
        if len(
            returned_scores
        ) > 1
        else 0.0
    )

    return {
        "classification": label_to_classification.get(
            best_label,
            "NONE"
        ),
        "best_score": best_score,
        "margin": best_score - second_score,
        "scope_classification": scope_classification,
        "scope_score": scope_score,
        "scope_margin": scope_margin,
        "matched_section": best_label
    }


def classification_risk_value(label):

    normalized = str(label or "SAFE").upper()

    if normalized == "SENSITIVE":
        return 3

    if normalized == "MEDIUM":
        return 2

    return 1


def apply_organization_policy_result(
    current_label,
    current_confidence,
    current_rule_score,
    current_reason,
    policy_record,
    policy_result,
    mandatory_protection_triggered=False
):

    final_label = str(current_label or "SAFE").upper()
    final_confidence = float(current_confidence or 0)
    final_rule_score = int(current_rule_score or 0)
    final_reason = str(current_reason or "").strip()

    if not policy_record.get("is_active"):
        return (
            final_label,
            final_confidence,
            final_rule_score,
            final_reason,
            False
        )

    policy_classification = str(
        policy_result.get("classification", "NONE")
    ).upper()

    policy_scope_classification = str(
        policy_result.get(
            "scope_classification",
            "NONE"
        )
    ).upper()

    policy_score = float(
        policy_result.get("best_score", 0)
    )

    policy_margin = float(
        policy_result.get("margin", 0)
    )

    policy_scope_classification = str(
        policy_result.get(
            "scope_classification",
            "NONE"
        )
    ).upper()

    # Mandatory baseline detections cannot be weakened by a custom policy.
    # Examples include confirmed passwords, credentials, API keys, tokens,
    # payment-card data, PII, and concrete private medical information.
    if mandatory_protection_triggered:
        return (
            final_label,
            final_confidence,
            final_rule_score,
            final_reason,
            False
        )

    if (
        policy_scope_classification == "OUT_OF_SCOPE"
        and not mandatory_protection_triggered
        and not (
            policy_result.get("organization_domain")
            == policy_result.get("content_domain")
            and policy_classification in {
                "MEDIUM",
                "SENSITIVE"
            }
        )
    ):
        policy_classification = "OUT_OF_SCOPE"

    # Organization custom business policies apply only to content
    # inside the organization's scope. If the scope analysis already
    # determined that the content belongs to another domain, do not allow
    # generic MEDIUM/SENSITIVE policy matches to override that decision.
    if (
        policy_scope_classification == "OUT_OF_SCOPE"
        and not mandatory_protection_triggered
    ):
        policy_classification = "OUT_OF_SCOPE"

    if policy_classification == "OUT_OF_SCOPE":

        # Prevent a false scope rejection from deleting a stronger baseline
        # classification. Organization scope is an additional control layer,
        # not a replacement for the specialized business classifiers.
        if final_label in {"MEDIUM", "SENSITIVE"}:
            return (
                final_label,
                final_confidence,
                final_rule_score,
                final_reason,
                False
            )

        organization_name = (
            policy_record.get("organization_name")
            or "the organization"
        )

        final_label = "SAFE"
        final_confidence = max(
            final_confidence,
            policy_score
        )
        final_rule_score = 0
        final_reason = (
            f"The content belongs to a different industry or domain and is "
            f"outside {organization_name}'s active protected business scope. "
            "No mandatory sensitive information was detected, so the "
            "organization-specific business classification is SAFE."
        )

        return (
            final_label,
            final_confidence,
            final_rule_score,
            final_reason,
            True
        )

    # Do not let the custom policy downgrade a strong in-scope baseline
    # business classification from SENSITIVE to MEDIUM. Cross-domain
    # content is already converted to OUT_OF_SCOPE before this function.
    #
    # Example:
    # - Kitchen policy + reproducible restaurant recipe:
    #   baseline SENSITIVE must remain SENSITIVE.
    # - Hospital policy + restaurant recipe:
    #   scope enforcement changes the policy result to OUT_OF_SCOPE,
    #   so it still becomes SAFE.
    if (
        final_label == "SENSITIVE"
        and policy_classification == "MEDIUM"
        and str(
            policy_result.get(
                "scope_classification",
                "NONE"
            )
        ).upper() != "OUT_OF_SCOPE"
    ):
        organization_name = (
            policy_record.get("organization_name")
            or "the organization"
        )

        final_confidence = max(
            final_confidence,
            policy_score
        )

        final_rule_score = max(
            final_rule_score,
            88
        )

        final_reason = (
            f"The active policy for {organization_name} considered the content "
            "MEDIUM, but the established in-scope business classifier detected "
            "strong SENSITIVE evidence. The system preserves the higher-risk "
            "classification so a strongly confidential in-scope document is not "
            "downgraded before administrator review."
        )

        return (
            "SENSITIVE",
            final_confidence,
            final_rule_score,
            final_reason,
            True
        )

    required_score = {
        "SENSITIVE": ORGANIZATION_POLICY_SENSITIVE_THRESHOLD,
        "MEDIUM": ORGANIZATION_POLICY_MEDIUM_THRESHOLD,
        "SAFE": ORGANIZATION_POLICY_SAFE_THRESHOLD
    }.get(policy_classification)

    if (
        required_score is None
        or policy_score < required_score
        or policy_margin < ORGANIZATION_POLICY_MIN_MARGIN
    ):
        return (
            final_label,
            final_confidence,
            final_rule_score,
            final_reason,
            False
        )

    previous_label = final_label

    # For organization-controlled business information, the active policy
    # may classify content as SAFE, MEDIUM, or SENSITIVE.
    final_label = policy_classification
    final_confidence = max(
        final_confidence,
        policy_score
    )

    if final_label == "SENSITIVE":
        final_rule_score = max(
            final_rule_score,
            88
        )
    elif final_label == "MEDIUM":
        final_rule_score = 48
    else:
        final_rule_score = 0

    organization_name = (
        policy_record.get("organization_name")
        or "the organization"
    )

    final_reason = (
        f"The active custom policy for {organization_name} classified this "
        f"content as {final_label}. The policy may override customizable "
        "business classifications such as recipes, internal procedures, "
        "supplier information, software intellectual property, manufacturing "
        "formulas, and other organization-defined information. Mandatory "
        "protections for passwords, credentials, API keys, access tokens, "
        "private keys, payment data, PII, and concrete private medical records "
        "cannot be reduced by the custom policy."
    )

    if previous_label != final_label:
        final_reason += (
            f" The baseline classification was {previous_label}."
        )

    return (
        final_label,
        final_confidence,
        final_rule_score,
        final_reason,
        True
    )

def build_explanation(
    final_label,
    final_confidence,
    rule_score,
    reason,
    file_extension,
    secret_refinement_label="",
    secret_refinement_score=0,
    secret_refinement_margin=0,
    generic_label="",
    generic_score=0,
    business_label="",
    business_score=0,
    business_safe_score=0,
    business_margin=0,
    healthcare_refinement_label="",
    healthcare_refinement_score=0,
    healthcare_refinement_margin=0,
    restaurant_refinement_label="",
    restaurant_refinement_score=0,
    restaurant_refinement_margin=0,
    secret_structural_evidence=False,
    payment_card_structure=False,
    assignment_value_shape=False,
    long_token_like_shape=False,
    reproducible_formula_detail=False,
    organization_name="",
    organization_policy_version=0,
    organization_policy_classification="NONE",
    organization_policy_score=0,
    organization_policy_margin=0,
    organization_policy_applied=False,
    organization_scope_classification="NONE",
    organization_scope_score=0,
    organization_scope_margin=0
):

    final_score = int(final_confidence * 100)

    return {
        "ml_prediction": final_label,
        "ml_confidence": round(final_confidence, 2),
        "ml_score": final_score,
        "rule_score": rule_score,
        "reason": reason,
        "secret_refinement_label": secret_refinement_label,
        "secret_refinement_score": round(secret_refinement_score, 2),
        "secret_refinement_margin": round(secret_refinement_margin, 2),
        "generic_ai_label": generic_label,
        "generic_ai_score": round(generic_score, 2),
        "business_ai_label": business_label,
        "business_ai_score": round(business_score, 2),
        "business_safe_score": round(business_safe_score, 2),
        "business_margin_over_safe": round(business_margin, 2),
        "healthcare_refinement_label": healthcare_refinement_label,
        "healthcare_refinement_score": round(healthcare_refinement_score, 2),
        "healthcare_refinement_margin": round(healthcare_refinement_margin, 2),
        "restaurant_refinement_label": restaurant_refinement_label,
        "restaurant_refinement_score": round(restaurant_refinement_score, 2),
        "restaurant_refinement_margin": round(restaurant_refinement_margin, 2),
        "secret_structural_evidence": secret_structural_evidence,
        "payment_card_structure": payment_card_structure,
        "assignment_value_shape": assignment_value_shape,
        "long_token_like_shape": long_token_like_shape,
        "reproducible_formula_detail": reproducible_formula_detail,
        "organization_name": organization_name,
        "organization_policy_version": organization_policy_version,
        "organization_policy_classification": organization_policy_classification,
        "organization_policy_score": round(organization_policy_score, 2),
        "organization_policy_margin": round(organization_policy_margin, 2),
        "organization_policy_applied": organization_policy_applied,
        "organization_scope_classification": organization_scope_classification,
        "organization_scope_score": round(organization_scope_score, 2),
        "organization_scope_margin": round(organization_scope_margin, 2),
        "file_extension": file_extension,
        "extraction_status": "TEXT_EXTRACTED"
    }


def run_payment_number_ai_analysis(text):
    """
    AI semantic refinement for card-shaped numbers.

    This does not use a predefined keyword list.
    The rule only detects that a 13 to 19 digit card-like structure exists.
    The AI model decides whether the surrounding context means exposed
    payment-card / financial-card data or just an unrelated number.
    """

    candidate_text = text[:3000]

    labels = [
        "the text exposes a real payment card or financial card number",
        "the text contains an unrelated long number, code, chapter number, example, or identifier",
        "the text is harmless general content"
    ]

    try:
        result = classifier(
            candidate_text,
            labels,
            multi_label=False
        )

        top_label = result["labels"][0]
        top_score = float(result["scores"][0])

        second_score = 0.0

        if len(result["scores"]) > 1:
            second_score = float(result["scores"][1])

        margin = top_score - second_score

        return top_label, top_score, margin

    except Exception:
        return "", 0.0, 0.0


def get_baseline_business_domain(
    business_label,
    business_score,
    business_margin
):

    if (
        business_label == HEALTHCARE_LABEL
        and business_score >= 0.55
    ):
        return "HEALTHCARE"

    if (
        business_label == RESTAURANT_LABEL
        and business_score >= 0.55
    ):
        return "FOOD_HOSPITALITY"

    if (
        business_label == SOFTWARE_LABEL
        and business_score >= 0.55
    ):
        return "SOFTWARE_TECHNOLOGY"

    if (
        business_label == CLOTHING_LABEL
        and business_score >= 0.55
    ):
        return "RETAIL_COMMERCE"

    return "NONE"


def enforce_active_organization_scope(
    policy_result,
    policy_record,
    business_label,
    business_score,
    business_margin,
    mandatory_protection_triggered
):

    result = dict(
        policy_result or {}
    )

    if (
        not policy_record.get(
            "is_active"
        )
        or mandatory_protection_triggered
    ):
        return result

    policy_sections = parse_organization_policy_sections(
        policy_record.get(
            "policy_prompt",
            ""
        )
    )

    scope_text = join_policy_section(
        policy_sections.get(
            "scope",
            []
        )
    )

    organization_domain_result = run_organization_domain_analysis(
        scope_text
    )

    organization_domain = organization_domain_result.get(
        "domain",
        "NONE"
    )

    baseline_content_domain = get_baseline_business_domain(
        business_label,
        business_score,
        business_margin
    )

    recognized_domains = set(
        ORGANIZATION_DOMAIN_LABELS.keys()
    ) - {
        "GENERAL_OTHER"
    }

    if (
        organization_domain in recognized_domains
        and baseline_content_domain in recognized_domains
        and organization_domain != baseline_content_domain
    ):
        organization_score = float(
            organization_domain_result.get(
                "score",
                0
            )
        )

        organization_margin = float(
            organization_domain_result.get(
                "margin",
                0
            )
        )

        result.update(
            {
                "classification": "OUT_OF_SCOPE",
                "best_score": max(
                    organization_score,
                    float(
                        business_score or 0
                    )
                ),
                "margin": max(
                    organization_margin,
                    float(
                        business_margin or 0
                    )
                ),
                "scope_classification": "OUT_OF_SCOPE",
                "scope_score": max(
                    organization_score,
                    float(
                        business_score or 0
                    )
                ),
                "scope_margin": max(
                    organization_margin,
                    float(
                        business_margin or 0
                    )
                ),
                "matched_section": (
                    "The active organization scope belongs to "
                    f"{organization_domain}, while the established baseline "
                    "business classifier identified the content as "
                    f"{baseline_content_domain}."
                )
            }
        )

    return result



def enforce_in_scope_business_classification(
    policy_result,
    policy_record,
    business_label,
    business_score,
    business_margin,
    restaurant_refinement_label,
    restaurant_refinement_score,
    restaurant_refinement_margin,
    reproducible_formula_detail,
    healthcare_identity_context,
    concrete_clinical_detail,
    incomplete_patient_file_context,
    mandatory_protection_triggered
):

    result = dict(
        policy_result or {}
    )

    if (
        not policy_record.get(
            "is_active"
        )
        or mandatory_protection_triggered
    ):
        return result

    policy_sections = parse_organization_policy_sections(
        policy_record.get(
            "policy_prompt",
            ""
        )
    )

    scope_text = join_policy_section(
        policy_sections.get(
            "scope",
            []
        )
    )

    organization_domain_result = run_organization_domain_analysis(
        scope_text
    )

    organization_domain = organization_domain_result.get(
        "domain",
        "NONE"
    )

    baseline_content_domain = get_baseline_business_domain(
        business_label,
        business_score,
        business_margin
    )

    # When the active organization and the content belong to the same
    # recognized business domain, strong baseline evidence must not be
    # weakened to MEDIUM by the custom policy model.
    if (
        organization_domain == "FOOD_HOSPITALITY"
        and baseline_content_domain == "FOOD_HOSPITALITY"
        and business_label == RESTAURANT_LABEL
        and reproducible_formula_detail
    ):
        strong_restaurant_evidence = (
            restaurant_refinement_label
            == RESTAURANT_REFINEMENT_SENSITIVE_LABEL
            and restaurant_refinement_score
            >= RESTAURANT_REFINEMENT_SENSITIVE_THRESHOLD
            and restaurant_refinement_margin
            >= RESTAURANT_REFINEMENT_MIN_MARGIN
        )

        if (
            strong_restaurant_evidence
            or business_score >= BUSINESS_MEDIUM_THRESHOLD
        ):
            result.update(
                {
                    "classification": "SENSITIVE",
                    "best_score": max(
                        0.99,
                        float(
                            business_score or 0
                        ),
                        float(
                            restaurant_refinement_score or 0
                        )
                    ),
                    "margin": max(
                        0.20,
                        float(
                            business_margin or 0
                        ),
                        float(
                            restaurant_refinement_margin or 0
                        )
                    ),
                    "scope_classification": "IN_SCOPE",
                    "scope_score": max(
                        float(
                            organization_domain_result.get(
                                "score",
                                0
                            )
                        ),
                        float(
                            business_score or 0
                        )
                    ),
                    "scope_margin": max(
                        float(
                            organization_domain_result.get(
                                "margin",
                                0
                            )
                        ),
                        float(
                            business_margin or 0
                        )
                    ),
                    "matched_section": (
                        "The active organization and the document both belong "
                        "to the food-service domain. The document contains a "
                        "reproducible internal recipe with concrete quantities, "
                        "preparation steps, timing, or storage instructions."
                    )
                }
            )

            return result

    if (
        organization_domain == "SOFTWARE_TECHNOLOGY"
        and baseline_content_domain == "SOFTWARE_TECHNOLOGY"
        and business_label == SOFTWARE_LABEL
        and business_score >= BUSINESS_SENSITIVE_THRESHOLD
        and business_margin >= BUSINESS_MIN_MARGIN
    ):
        result.update(
            {
                "classification": "SENSITIVE",
                "best_score": max(
                    0.95,
                    float(
                        business_score or 0
                    )
                ),
                "margin": max(
                    0.15,
                    float(
                        business_margin or 0
                    )
                ),
                "scope_classification": "IN_SCOPE",
                "scope_score": max(
                    float(
                        organization_domain_result.get(
                            "score",
                            0
                        )
                    ),
                    float(
                        business_score or 0
                    )
                ),
                "scope_margin": max(
                    float(
                        organization_domain_result.get(
                            "margin",
                            0
                        )
                    ),
                    float(
                        business_margin or 0
                    )
                ),
                "matched_section": (
                    "The active organization and the document both belong "
                    "to the software-technology domain, and strong proprietary "
                    "software evidence was detected."
                )
            }
        )

        return result

    if (
        organization_domain == "RETAIL_COMMERCE"
        and baseline_content_domain == "RETAIL_COMMERCE"
        and business_label == CLOTHING_LABEL
    ):
        if (
            business_score >= BUSINESS_SENSITIVE_THRESHOLD
            and business_margin >= BUSINESS_MIN_MARGIN
        ):
            result.update(
                {
                    "classification": "SENSITIVE",
                    "best_score": max(
                        0.95,
                        float(business_score or 0)
                    ),
                    "margin": max(
                        0.15,
                        float(business_margin or 0)
                    ),
                    "scope_classification": "IN_SCOPE",
                    "scope_score": max(
                        float(
                            organization_domain_result.get(
                                "score",
                                0
                            )
                        ),
                        float(business_score or 0)
                    ),
                    "scope_margin": max(
                        float(
                            organization_domain_result.get(
                                "margin",
                                0
                            )
                        ),
                        float(business_margin or 0)
                    ),
                    "matched_section": (
                        "The active organization and the document both belong "
                        "to the retail/clothing domain, and confidential clothing "
                        "business information was detected."
                    )
                }
            )

            return result

        elif (
            business_score >= BUSINESS_MEDIUM_THRESHOLD
            and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
        ):
            result.update(
                {
                    "classification": "MEDIUM",
                    "best_score": max(
                        0.90,
                        float(business_score or 0)
                    ),
                    "margin": max(
                        0.12,
                        float(business_margin or 0)
                    ),
                    "scope_classification": "IN_SCOPE",
                    "scope_score": max(
                        float(
                            organization_domain_result.get(
                                "score",
                                0
                            )
                        ),
                        float(business_score or 0)
                    ),
                    "scope_margin": max(
                        float(
                            organization_domain_result.get(
                                "margin",
                                0
                            )
                        ),
                        float(business_margin or 0)
                    ),
                    "matched_section": (
                        "The active organization and the document both belong "
                        "to the retail/clothing domain, and possible confidential "
                        "retail information requires administrator review."
                    )
                }
            )

            return result

    if (
        organization_domain == "HEALTHCARE"
        and incomplete_patient_file_context
        and not concrete_clinical_detail
    ):
        result.update(
            {
                "classification": "MEDIUM",
                "best_score": max(
                    0.90,
                    float(
                        business_score or 0
                    )
                ),
                "margin": max(
                    0.15,
                    float(
                        business_margin or 0
                    )
                ),
                "scope_classification": "IN_SCOPE",
                "scope_score": max(
                    float(
                        organization_domain_result.get(
                            "score",
                            0
                        )
                    ),
                    float(
                        business_score or 0
                    )
                ),
                "scope_margin": max(
                    float(
                        organization_domain_result.get(
                            "margin",
                            0
                        )
                    ),
                    float(
                        business_margin or 0
                    )
                ),
                "matched_section": (
                    "The active organization is healthcare, and the document "
                    "contains an incomplete patient-file reference without "
                    "concrete patient identity or clinical details. It requires "
                    "administrator review."
                )
            }
        )

        return result

    if (
        organization_domain == "HEALTHCARE"
        and baseline_content_domain == "HEALTHCARE"
        and healthcare_identity_context
        and concrete_clinical_detail
    ):
        result.update(
            {
                "classification": "SENSITIVE",
                "best_score": max(
                    0.99,
                    float(
                        business_score or 0
                    )
                ),
                "margin": max(
                    0.20,
                    float(
                        business_margin or 0
                    )
                ),
                "scope_classification": "IN_SCOPE",
                "scope_score": max(
                    float(
                        organization_domain_result.get(
                            "score",
                            0
                        )
                    ),
                    float(
                        business_score or 0
                    )
                ),
                "scope_margin": max(
                    float(
                        organization_domain_result.get(
                            "margin",
                            0
                        )
                    ),
                    float(
                        business_margin or 0
                    )
                ),
                "matched_section": (
                    "The active organization and the document both belong "
                    "to the healthcare domain, and the document contains "
                    "patient identity together with concrete clinical details."
                )
            }
        )

    return result



def predict_file(file_path):

    try:
        content = read_file_content(file_path)
    except Exception:
        content = ""

    file_extension = os.path.splitext(file_path)[1].lower()

    if not content.strip():

        explanation = {
            "ml_prediction": "SAFE",
            "ml_confidence": 0,
            "ml_score": 0,
            "rule_score": 0,
            "reason": (
                "No readable text content was extracted from the file. "
                "For images or scanned PDFs, verify that Tesseract OCR, "
                "Pillow, pytesseract, and PyMuPDF are installed and that "
                "the image contains clear readable text."
            ),
            "file_extension": file_extension,
            "extraction_status": "NO_TEXT_EXTRACTED"
        }

        return "SAFE", 0, [], "", explanation

    secret_refinement = run_secret_refinement_ai_analysis(content)
    generic_result = run_generic_ai_analysis(content)
    business_result = run_business_ai_analysis(content)

    secret_label = secret_refinement["best_label"]
    secret_score = secret_refinement["best_score"]
    secret_margin = secret_refinement["margin"]

    generic_label = generic_result["best_label"]
    generic_score = generic_result["best_score"]

    business_label = business_result["best_label"]
    business_score = business_result["best_score"]
    business_safe_score = business_result["safe_score"]
    business_margin = business_result["margin_over_safe"]

    healthcare_refinement_label = ""
    healthcare_refinement_score = 0
    healthcare_refinement_margin = 0

    restaurant_refinement_label = ""
    restaurant_refinement_score = 0
    restaurant_refinement_margin = 0

    payment_card_structure = has_payment_card_structure(content)
    assignment_value_shape = has_assignment_value_shape(content)
    long_token_like_shape = has_long_token_like_shape(content)
    secret_structural_evidence = has_secret_structural_evidence(content)
    credential_exposure_shape = has_credential_exposure_shape(content)
    weak_credential_context = has_weak_credential_context(content)

    reproducible_formula_detail = has_reproducible_formula_detail(content)

    payment_number_ai_label = ""
    payment_number_ai_score = 0.0
    payment_number_ai_margin = 0.0

    if has_valid_payment_card_number(
        content
    ):

        masked_content = mask_sensitive_content(
            content
        )

        explanation = {
            "ml_prediction": "SENSITIVE",
            "ml_confidence": 0.99,
            "ml_score": 99,
            "rule_score": 95,
            "reason": (
                "The content contains a complete payment-card number "
                "that passed structural card-number validation. Full "
                "payment-card numbers are protected by the mandatory "
                "DLP baseline regardless of the active organization "
                "policy. Sensitive values are masked below ->\n"
                f"{masked_content}"
            ),
            "payment_number_ai_label": (
                "validated complete payment-card number"
            ),
            "payment_number_ai_score": 0.99,
            "payment_number_ai_margin": 0.99,
            "structural_payment_card_number": True,
            "mandatory_protection_triggered": True,
            "file_extension": file_extension,
            "extraction_status": "OK"
        }

        return (
            "SENSITIVE",
            explanation["ml_score"],
            [],
            content,
            explanation
        )

    if has_card_number_shape(content):
        (
            payment_number_ai_label,
            payment_number_ai_score,
            payment_number_ai_margin
        ) = run_payment_number_ai_analysis(content)

        if (
            payment_number_ai_label
            == "the text exposes a real payment card or financial card number"
            and payment_number_ai_score >= 0.55
            and payment_number_ai_margin >= 0.00
        ):
            masked_content = mask_sensitive_content(content)

            explanation = {
                "ml_prediction": "SENSITIVE",
                "ml_confidence": round(payment_number_ai_score, 2),
                "ml_score": 95,
                "rule_score": 0,
                "reason": (
                    "A card-like numeric structure was found, and AI semantic "
                    "analysis determined that the surrounding context describes "
                    "exposed payment-card or financial-card data. Sensitive "
                    "values are masked below ->\n"
                    f"{masked_content}"
                ),
                "payment_number_ai_label": payment_number_ai_label,
                "payment_number_ai_score": payment_number_ai_score,
                "payment_number_ai_margin": payment_number_ai_margin,
                "structural_payment_card_number": True,
                "file_extension": file_extension,
                "extraction_status": "OK"
            }

            return "SENSITIVE", explanation["ml_score"], [], content, explanation

    medium_risk_context = has_medium_risk_context(content)
    healthcare_identity_context = has_healthcare_identity_context(content)
    concrete_clinical_detail = has_concrete_clinical_detail(content)
    public_healthcare_information = has_public_healthcare_information(
        content
    )
    incomplete_patient_file_context = has_incomplete_patient_file_context(
        content
    )

    if credential_exposure_shape:

        explanation = build_explanation(
            "SENSITIVE",
            0.99,
            95,
            (
                "A credential exposure pattern was detected. The content contains "
                "a password, passcode, API key, secret key, access token, auth token, "
                "or private key assignment with an exposed value."
            ),
            file_extension,
            secret_label,
            secret_score,
            secret_margin,
            generic_label,
            generic_score,
            business_label,
            business_score,
            business_safe_score,
            business_margin,
            healthcare_refinement_label,
            healthcare_refinement_score,
            healthcare_refinement_margin,
            restaurant_refinement_label,
            restaurant_refinement_score,
            restaurant_refinement_margin,
            True,
            payment_card_structure,
            assignment_value_shape,
            long_token_like_shape,
            reproducible_formula_detail
        )

        return "SENSITIVE", explanation["ml_score"], [], content, explanation

    if weak_credential_context:

        explanation = build_explanation(
            "MEDIUM",
            0.70,
            50,
            (
                "A possible credential or token-like configuration value was detected, "
                "but it is not strong enough to confirm a real usable secret. The file "
                "requires dashboard review."
            ),
            file_extension,
            secret_label,
            secret_score,
            secret_margin,
            generic_label,
            generic_score,
            business_label,
            business_score,
            business_safe_score,
            business_margin,
            healthcare_refinement_label,
            healthcare_refinement_score,
            healthcare_refinement_margin,
            restaurant_refinement_label,
            restaurant_refinement_score,
            restaurant_refinement_margin,
            False,
            payment_card_structure,
            assignment_value_shape,
            long_token_like_shape,
            reproducible_formula_detail
        )

        return "MEDIUM", explanation["ml_score"], [], content, explanation

    business_sensitive_labels = [
        HEALTHCARE_LABEL,
        RESTAURANT_LABEL,
        SOFTWARE_LABEL
    ]

    if (
        secret_label == SECRET_EXPOSED_LABEL
        and secret_score >= SECRET_SENSITIVE_THRESHOLD
        and secret_structural_evidence
        and not (
            business_label == HEALTHCARE_LABEL
            and business_score >= BUSINESS_MEDIUM_THRESHOLD
            and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
        )
    ):

        final_confidence = max(secret_score, 0.62)
        masked_content = mask_sensitive_content(content)

        explanation = build_explanation(
            "SENSITIVE",
            final_confidence,
            90,
            (
                "AI semantic analysis determined that the text exposes actual "
                "usable sensitive data, and structural evidence was found in "
                "the content, such as a payment-card pattern, credential/key "
                "assignment pattern, or token-like value shape. Sensitive values "
                "are masked below ->\n"
                f"{masked_content}"
            ),
            file_extension,
            secret_label,
            secret_score,
            secret_margin,
            generic_label,
            generic_score,
            business_label,
            business_score,
            business_safe_score,
            business_margin,
            healthcare_refinement_label,
            healthcare_refinement_score,
            healthcare_refinement_margin,
            restaurant_refinement_label,
            restaurant_refinement_score,
            restaurant_refinement_margin,
            secret_structural_evidence,
            payment_card_structure,
            assignment_value_shape,
            long_token_like_shape,
            reproducible_formula_detail
        )

        return "SENSITIVE", explanation["ml_score"], [], content, explanation

    if (
        secret_label == SECRET_EXPOSED_LABEL
        and secret_score >= SECRET_MEDIUM_THRESHOLD
        and secret_structural_evidence
        and not (
            business_label == HEALTHCARE_LABEL
            and business_score >= BUSINESS_MEDIUM_THRESHOLD
            and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
        )
    ):

        final_confidence = max(secret_score, 0.55)

        explanation = build_explanation(
            "MEDIUM",
            final_confidence,
            50,
            (
                "AI semantic analysis found possible exposed sensitive data and "
                "structural evidence exists, but confidence was not strong enough "
                "for automatic blocking."
            ),
            file_extension,
            secret_label,
            secret_score,
            secret_margin,
            generic_label,
            generic_score,
            business_label,
            business_score,
            business_safe_score,
            business_margin,
            healthcare_refinement_label,
            healthcare_refinement_score,
            healthcare_refinement_margin,
            restaurant_refinement_label,
            restaurant_refinement_score,
            restaurant_refinement_margin,
            secret_structural_evidence,
            payment_card_structure,
            assignment_value_shape,
            long_token_like_shape,
            reproducible_formula_detail
        )

        return "MEDIUM", explanation["ml_score"], [], content, explanation

    final_label = "SAFE"
    final_confidence = max(secret_score, generic_score, business_score)
    rule_score = 0
    reason_text = build_safe_content_reason(content)

    if healthcare_identity_context and concrete_clinical_detail:

        final_label = "SENSITIVE"
        final_confidence = max(final_confidence, 0.92)
        rule_score = 90
        reason_text = (
            "Private healthcare information was detected. The content contains "
            "patient or medical-record context together with concrete clinical "
            "details such as diagnosis, prescription, treatment, laboratory result, "
            "medical condition, or clinical notes."
        )

    elif healthcare_identity_context:

        final_label = "MEDIUM"
        final_confidence = max(final_confidence, 0.72)
        rule_score = 45
        reason_text = (
            "Healthcare-related identity or patient-file context was detected, "
            "but no concrete diagnosis, prescription, lab result, treatment, "
            "medical condition, or clinical detail was confirmed. The file "
            "requires dashboard review."
        )

    elif (
        business_label == HEALTHCARE_LABEL
        and has_explicit_healthcare_context(content)
    ):

        healthcare_result = run_healthcare_refinement_ai_analysis(content)

        healthcare_refinement_label = healthcare_result["best_label"]
        healthcare_refinement_score = healthcare_result["best_score"]
        healthcare_refinement_margin = healthcare_result["margin"]

        final_confidence = max(final_confidence, healthcare_refinement_score)

        if (
            healthcare_refinement_label == HEALTHCARE_REFINEMENT_SENSITIVE_LABEL
            and healthcare_refinement_score >= HEALTHCARE_REFINEMENT_MEDIUM_THRESHOLD
            and healthcare_refinement_margin >= HEALTHCARE_REFINEMENT_MEDIUM_MARGIN
            and concrete_clinical_detail
            and count_words(content) >= 8
        ):

            final_label = "SENSITIVE"
            rule_score = 90
            reason_text = (
                "Focused AI healthcare analysis detected concrete private patient "
                "clinical information, such as diagnosis, prescription, treatment, "
                "lab result, medical condition, or clinical details."
            )

        else:

            final_label = "SAFE"
            rule_score = 0
            reason_text = build_safe_content_reason(content)

    elif business_label == RESTAURANT_LABEL:

        restaurant_result = run_restaurant_refinement_ai_analysis(content)

        restaurant_refinement_label = restaurant_result["best_label"]
        restaurant_refinement_score = restaurant_result["best_score"]
        restaurant_refinement_margin = restaurant_result["margin"]

        final_confidence = max(final_confidence, restaurant_refinement_score)

        if reproducible_formula_detail:

            if (
                restaurant_refinement_label == RESTAURANT_REFINEMENT_SENSITIVE_LABEL
                and restaurant_refinement_score >= RESTAURANT_REFINEMENT_SENSITIVE_THRESHOLD
                and restaurant_refinement_margin >= RESTAURANT_REFINEMENT_MIN_MARGIN
            ):

                final_label = "SENSITIVE"
                rule_score = 92
                reason_text = (
                    "Focused AI restaurant analysis detected business-sensitive "
                    "recipe or food production information. The content contains "
                    "enough concrete formula or preparation details to reproduce "
                    "an internal restaurant menu item."
                )

            elif (
                business_score >= BUSINESS_MEDIUM_THRESHOLD
                and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
            ):

                final_label = "MEDIUM"
                rule_score = 50
                reason_text = (
                    "AI detected possible restaurant food-production information. "
                    "The text contains some reproducible formula structure and "
                    "requires review."
                )

        elif has_medium_restaurant_context(content):

            final_label = "MEDIUM"
            rule_score = 45
            reason_text = (
                "Possible internal restaurant food-production information was detected, "
                "but the content does not contain enough exact formula details for "
                "automatic blocking. The file requires dashboard review."
            )

        else:

            if (
                restaurant_refinement_label in [
                    RESTAURANT_REFINEMENT_PUBLIC_LABEL,
                    RESTAURANT_REFINEMENT_OPERATIONAL_LABEL
                ]
                and restaurant_refinement_score >= RESTAURANT_REFINEMENT_SAFE_THRESHOLD
            ):

                final_label = "SAFE"
                rule_score = 0
                reason_text = (
                    "Focused AI restaurant analysis determined that the text does "
                    "not reveal enough concrete formula or food-production details "
                    "to reproduce an internal menu item."
                )

    elif business_label == CLOTHING_LABEL:

        if (
            business_score >= BUSINESS_SENSITIVE_THRESHOLD
            and business_margin >= BUSINESS_MIN_MARGIN
        ):

            final_label = "SENSITIVE"
            rule_score = 90
            reason_text = get_business_reason(business_label)

        elif (
            business_score >= BUSINESS_MEDIUM_THRESHOLD
            and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
        ):

            final_label = "MEDIUM"
            rule_score = 50
            reason_text = (
                get_business_reason(business_label)
                + " Confidence was not high enough for automatic blocking, "
                + "so the content is marked for review."
            )

    elif business_label == SOFTWARE_LABEL:

        if (
            business_score >= BUSINESS_SENSITIVE_THRESHOLD
            and business_margin >= BUSINESS_MIN_MARGIN
        ):

            final_label = "SENSITIVE"
            rule_score = 90
            reason_text = get_business_reason(business_label)

        elif (
            business_score >= BUSINESS_MEDIUM_THRESHOLD
            and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
        ):

            final_label = "MEDIUM"
            rule_score = 50
            reason_text = (
                get_business_reason(business_label)
                + " Confidence was not high enough for automatic blocking, "
                + "so the content is marked for review."
            )

    if final_label == "SAFE" and has_medium_restaurant_context(content):

        final_label = "MEDIUM"
        final_confidence = max(final_confidence, 0.70)
        rule_score = 45
        reason_text = (
            "Possible internal restaurant food-production information was detected, "
            "but the content does not contain enough exact formula details for "
            "automatic blocking. The file requires dashboard review."
        )

    if final_label == "SAFE" and has_medium_software_context(content):

        final_label = "MEDIUM"
        final_confidence = max(final_confidence, 0.70)
        rule_score = 45
        reason_text = (
            "Possible internal software or system-design information was detected, "
            "but the content does not expose enough concrete proprietary implementation "
            "detail for automatic blocking. The file requires dashboard review."
        )

    if final_label == "SAFE" and has_confidential_warning_context(content):

        final_label = "MEDIUM"
        final_confidence = max(final_confidence, 0.68)
        rule_score = 40
        reason_text = (
            "The file contains a confidentiality, sensitivity, restricted, private, "
            "or do-not-share warning with document or business context. No concrete "
            "secret was exposed, so it requires dashboard review instead of blocking."
        )

    if final_label == "SAFE":

        if (
            generic_label == GENERIC_SENSITIVE_LABEL
            and generic_score >= GENERIC_SENSITIVE_THRESHOLD
            and count_words(content) > 8
            and (
                secret_structural_evidence
                or (
                    business_label in business_sensitive_labels
                    and business_score >= BUSINESS_SENSITIVE_THRESHOLD
                    and business_margin >= BUSINESS_MIN_MARGIN
                )
            )
        ):

            final_label = "SENSITIVE"
            rule_score = 85
            reason_text = (
                "AI semantic analysis detected highly confident confidential "
                "or sensitive information, supported by sensitive context or "
                "structural evidence."
            )

        elif (
            generic_label == GENERIC_SENSITIVE_LABEL
            and generic_score >= GENERIC_MEDIUM_THRESHOLD
            and count_words(content) > 8
            and (
                secret_structural_evidence
                or (
                    business_label in business_sensitive_labels
                    and business_score >= BUSINESS_MEDIUM_THRESHOLD
                    and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
                )
            )
        ):

            final_label = "MEDIUM"
            rule_score = 40
            reason_text = (
                "AI semantic analysis detected possible confidential or sensitive "
                "information, supported by sensitive context or structural evidence."
            )

    if (
        final_label in ["MEDIUM", "SENSITIVE"]
        and not secret_structural_evidence
        and not reproducible_formula_detail
        and not medium_risk_context
        and not weak_credential_context
        and not healthcare_identity_context
        and not (
            business_label == HEALTHCARE_LABEL
            and healthcare_refinement_label == HEALTHCARE_REFINEMENT_SENSITIVE_LABEL
            and concrete_clinical_detail
            and healthcare_refinement_score >= HEALTHCARE_REFINEMENT_MEDIUM_THRESHOLD
            and healthcare_refinement_margin >= HEALTHCARE_REFINEMENT_MEDIUM_MARGIN
            and count_words(content) >= 8
        )
        and not (
            business_label == SOFTWARE_LABEL
            and business_score >= BUSINESS_MEDIUM_THRESHOLD
            and business_margin >= BUSINESS_MEDIUM_MIN_MARGIN
        )
    ):

        final_label = "SAFE"
        rule_score = 0
        reason_text = build_safe_content_reason(content)

    organization_policy = get_active_organization_policy()

    organization_policy_result = (
        run_organization_policy_ai_analysis(
            content,
            organization_policy.get(
                "policy_prompt",
                ""
            ),
            organization_policy.get(
                "organization_name",
                ""
            )
        )
    )

    mandatory_protection_triggered = (
        credential_exposure_shape
        or weak_credential_context
        or payment_card_structure
        or has_valid_payment_card_number(
            content
        )
        or (
            healthcare_identity_context
            and concrete_clinical_detail
        )
    )

    organization_policy_result = enforce_active_organization_scope(
        organization_policy_result,
        organization_policy,
        business_label,
        business_score,
        business_margin,
        mandatory_protection_triggered
    )

    organization_policy_result = enforce_in_scope_business_classification(
        organization_policy_result,
        organization_policy,
        business_label,
        business_score,
        business_margin,
        restaurant_refinement_label,
        restaurant_refinement_score,
        restaurant_refinement_margin,
        reproducible_formula_detail,
        healthcare_identity_context,
        concrete_clinical_detail,
        incomplete_patient_file_context,
        mandatory_protection_triggered
    )

    (
        final_label,
        final_confidence,
        rule_score,
        reason_text,
        organization_policy_applied
    ) = apply_organization_policy_result(
        final_label,
        final_confidence,
        rule_score,
        reason_text,
        organization_policy,
        organization_policy_result,
        mandatory_protection_triggered=(
            mandatory_protection_triggered
        )
    )

    if (
        public_healthcare_information
        and not credential_exposure_shape
        and not weak_credential_context
        and not payment_card_structure
        and not has_valid_payment_card_number(
            content
        )
        and not healthcare_identity_context
        and not concrete_clinical_detail
    ):
        final_label = "SAFE"
        final_confidence = max(
            final_confidence,
            0.95
        )
        rule_score = 0
        reason_text = (
            "The content contains public hospital or clinic operational "
            "information, such as visiting hours, opening hours, contact "
            "details, directions, or visitor information. It does not expose "
            "a patient identity, diagnosis, prescription, treatment, lab "
            "result, medical history, credential, or payment-card data."
        )
        organization_policy_applied = True

    explanation = build_explanation(
        final_label,
        final_confidence,
        rule_score,
        reason_text,
        file_extension,
        secret_label,
        secret_score,
        secret_margin,
        generic_label,
        generic_score,
        business_label,
        business_score,
        business_safe_score,
        business_margin,
        healthcare_refinement_label,
        healthcare_refinement_score,
        healthcare_refinement_margin,
        restaurant_refinement_label,
        restaurant_refinement_score,
        restaurant_refinement_margin,
        secret_structural_evidence,
        payment_card_structure,
        assignment_value_shape,
        long_token_like_shape,
        reproducible_formula_detail,
        organization_policy.get(
            "organization_name",
            ""
        ),
        organization_policy.get(
            "policy_version",
            0
        ),
        organization_policy_result.get(
            "classification",
            "NONE"
        ),
        organization_policy_result.get(
            "best_score",
            0
        ),
        organization_policy_result.get(
            "margin",
            0
        ),
        organization_policy_applied,
        organization_policy_result.get(
            "scope_classification",
            "NONE"
        ),
        organization_policy_result.get(
            "scope_score",
            0
        ),
        organization_policy_result.get(
            "scope_margin",
            0
        )
    )

    return final_label, explanation["ml_score"], [], content, explanation
